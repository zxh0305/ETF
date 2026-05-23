#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ETF指数智能投资顾问 - 系统架构与数据报告 PPT
生成时间: 2026-05-23
"""

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json, os
from datetime import datetime

# ============================================================
# 亮色主题配色
# ============================================================
C_NAVY   = RGBColor(0x1A, 0x3A, 0x5C)   # 深蓝（标题）
C_BLUE   = RGBColor(0x44, 0x88, 0xDD)   # 亮蓝
C_SKY    = RGBColor(0x5B, 0xB0, 0xE8)   # 天蓝
C_GOLD   = RGBColor(0xF0, 0xA5, 0x00)   # 金黄
C_GRAY   = RGBColor(0x44, 0x44, 0x44)
C_LGRAY  = RGBColor(0x66, 0x66, 0x66)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)   # 绿色
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)   # 红色
C_ORANGE = RGBColor(0xE8, 0x7B, 0x19)
C_BG     = RGBColor(0xF5, 0xF8, 0xFF)   # 浅蓝白背景
C_CARD   = RGBColor(0xE8, 0xF4, 0xFD)   # 卡片浅蓝
C_YELLOW = RGBColor(0xFF, 0xF3, 0xCD)   # 黄色提示

prs = Presentation()
prs.slide_width  = Emu(9144000)
prs.slide_height = Emu(6858000)

TITLE_H   = 1000000
GOLD_BAR  = 60000
CONTENT_Y = TITLE_H + GOLD_BAR + 80000
MARGIN_L  = 300000
MARGIN_R  = 300000
CONTENT_W = 9144000 - MARGIN_L - MARGIN_R
GAP       = 70000

def R(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def TB(slide, l, t, w, h, text, size, bold=False, color=C_WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    """深色背景文字框（白字/亮字）"""
    s = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    s.text_frame.word_wrap = True
    for para_text in text.split('\n'):
        para = s.text_frame.add_paragraph() if s.text_frame.paragraphs[0].text or s.text_frame.paragraphs[0].runs else s.text_frame.paragraphs[0]
        if s.text_frame.paragraphs[0].text:
            para = s.text_frame.add_paragraph()
        else:
            para = s.text_frame.paragraphs[0]
        para.text = para_text
        para.alignment = align
        if size:
            for run in para.runs:
                run.font.size = size
        if bold:
            for run in para.runs:
                run.font.bold = True
        if italic:
            for run in para.runs:
                run.font.italic = True
        for run in para.runs:
            run.font.color.rgb = color
    return s

def T(slide, l, t, w, h, text, size, bold=False, color=C_LGRAY,
       align=PP_ALIGN.LEFT, italic=False):
    """浅色背景文字框（深色文字）"""
    s = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    s.text_frame.word_wrap = True
    lines = text.split('\n')
    for i, line_text in enumerate(lines):
        if i == 0:
            para = s.text_frame.paragraphs[0]
        else:
            para = s.text_frame.add_paragraph()
        para.text = line_text
        para.alignment = align
        if size:
            for run in para.runs:
                run.font.size = size
        if bold:
            for run in para.runs:
                run.font.bold = bold
        if italic:
            for run in para.runs:
                run.font.italic = italic
        for run in para.runs:
            run.font.color.rgb = color
    return s

def title_bar(slide, title_text, subtitle_text=""):
    """蓝色标题栏 + 金色底边"""
    R(slide, 0, 0, 9144000, TITLE_H, C_NAVY)
    R(slide, 0, TITLE_H, 9144000, GOLD_BAR, C_GOLD)
    TB(slide, MARGIN_L, 100000, CONTENT_W, 500000,
       title_text, Pt(28), bold=True, color=C_WHITE)
    if subtitle_text:
        TB(slide, MARGIN_L, 580000, CONTENT_W, 400000,
           subtitle_text, Pt(14), color=RGBColor(0xBB, 0xCC, 0xDD), italic=True)

def section_header(slide, text, color=C_BLUE):
    R(slide, MARGIN_L, CONTENT_Y, 300000, 280000, color)
    TB(slide, MARGIN_L + 350000, CONTENT_Y, CONTENT_W - 350000, 280000,
       text, Pt(16), bold=True, color=color)

def card(slide, l, t, w, h, title, lines, title_color=C_BLUE):
    """带标题的卡片"""
    R(slide, l, t, w, h, C_CARD)
    R(slide, l, t, w, 50000, title_color)
    TB(slide, l + 15000, t + 8000, w - 30000, 400000,
       title, Pt(13), bold=True, color=C_WHITE)
    TB(slide, l + 15000, t + 550000, w - 30000, h - 650000,
       '\n'.join(lines), Pt(11), color=C_GRAY)

# ============================================================
# 加载数据
# ============================================================
BASE_DIR = os.path.dirname(os.path.abspath(__file__)).replace('/scripts','')
OUTPUT_DIR = os.path.join(BASE_DIR, 'output')

with open(os.path.join(OUTPUT_DIR, 'low_valuation_candidates_latest.json')) as f:
    candidates = json.load(f)
with open(os.path.join(OUTPUT_DIR, 'portfolio_latest.json')) as f:
    portfolio = json.load(f)
with open(os.path.join(OUTPUT_DIR, 'daily_report_latest.txt'), encoding='utf-8') as f:
    daily_report = f.read()

formal_pool = candidates.get('formal_pool', [])
watch_list  = candidates.get('watch_list', [])
stats       = candidates.get('stats', {})
meta        = candidates.get('meta', {})
portfolio_items = portfolio.get('items', [])
portfolio_meta  = portfolio.get('meta', {})

# ============================================================
# PPT Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 全背景
R(slide, 0, 0, 9144000, 6858000, C_NAVY)

# 装饰条
R(slide, 0, 0, 200000, 6858000, C_BLUE)
R(slide, 200000, 0, 60000, 6858000, C_SKY)

# 金色分割线
R(slide, 500000, 3500000, 8144000, 60000, C_GOLD)
R(slide, 500000, 4500000, 8144000, 60000, C_GOLD)

# 主标题
TB(slide, 500000, 1200000, 8144000, 1500000,
   'ETF 指数智能投资顾问', Pt(52), bold=True, color=C_WHITE,
   align=PP_ALIGN.LEFT)

# 副标题
TB(slide, 500000, 2600000, 8144000, 700000,
   '系统架构 · 数据来源 · 筛选逻辑 · 今日报告',
   Pt(20), color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.LEFT)

# 底部信息
TB(slide, 500000, 5000000, 5000000, 500000,
   f'📅 报告日期: 2026-05-23\n📊 低估池: {stats.get("formal_pool",0)}只正式 + {stats.get("watch_list",0)}只关注  共{stats.get("total",0)}只A股 ETF',
   Pt(14), color=RGBColor(0x99, 0xBB, 0xDD), align=PP_ALIGN.LEFT)

# ============================================================
# PPT Slide 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, '内容概览', 'CONTENTS')

items = [
    ('01', '系统架构总览',      '全流程9步流水线，从数据采集到微信推送'),
    ('02', '四大数据源',         '乐咕乐股 · 申万行业 · 巨潮资讯 · 新浪行情'),
    ('03', '筛选条件详解',      'PE/PB分位 ≤30% · PEG<1 · 日均成交≥1亿'),
    ('04', '数据质量体系',      '✅真实数据 vs ⚠️估算数据，如何区分'),
    ('05', '为什么只有41只ETF', '宽基指数ETF为何普遍偏贵，申万行业分位逻辑'),
    ('06', '今日低估池',        '36只正式池 + 113只关注池完整列表'),
    ('07', '组合配置建议',      'Portfolio Engine 输出的推荐仓位（平衡模式）'),
    ('08', '风险提示',          '数据局限性与使用注意事项'),
]

for i, (num, title, desc) in enumerate(items):
    y = CONTENT_Y + 150000 + i * 620000
    R(slide, MARGIN_L, y, 700000, 520000, C_BLUE)
    TB(slide, MARGIN_L + 50000, y + 80000, 600000, 400000,
       num, Pt(22), bold=True, color=C_WHITE)
    TB(slide, MARGIN_L + 900000, y + 50000, 6000000, 400000,
       title, Pt(16), bold=True, color=C_NAVY)
    TB(slide, MARGIN_L + 900000, y + 380000, 7000000, 300000,
       desc, Pt(12), color=C_LGRAY, italic=True)

# ============================================================
# PPT Slide 3: 系统架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, '系统架构总览', '全天自动运行 · 3个Agent分工协作 · 每天09:20启动')

steps = [
    ('Step0', '申万行业\n估值刷新',    C_BLUE,   'akshare sw_index\n获取31个申万\n一级行业PE/PB'),
    ('Step1', '行情数据\n采集',         C_SKY,    '东方财富ETF列表\n全市场1400+只\n实时价格/成交额'),
    ('Step2', '估值补全',               C_BLUE,   '四数据源合并\nPE/PB分位计算\n申万行业穿透'),
    ('Step3', '低估筛选',               C_GOLD,   'PE%≤30%\nPEG<1\n日均≥1亿'),
    ('Step4', '变化对比',               C_GREEN,  '昨日vs今日\n新入/加仓\n减仓信号'),
    ('Step5', '日报生成',               C_ORANGE, 'Markdown日报\n微信推送模板\n数据质量标记'),
    ('Step6', '指数快照',               C_BLUE,   '每日存档\n宽基指数PE/PB\n分位历史'),
    ('Step7', 'Portfolio\nEngine',      C_RED,    '相关性分析\n风险风控\n仓位优化'),
]

step_w = (CONTENT_W - 7*GAP) // 8
for i, (sid, name, color, desc) in enumerate(steps):
    l = MARGIN_L + i*(step_w + GAP)
    R(slide, l, CONTENT_Y, step_w, 600000, color)
    TB(slide, l, CONTENT_Y + 80000, step_w, 450000,
       sid, Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    TB(slide, l + 10000, CONTENT_Y + 220000, step_w - 20000, 450000,
       name, Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # 箭头
    if i < 7:
        R(slide, l + step_w, CONTENT_Y + 200000, GAP, 100000, C_GOLD)
    # 描述
    TB(slide, l, CONTENT_Y + 700000, step_w, 900000,
       desc, Pt(9), color=C_GRAY, align=PP_ALIGN.CENTER)

# 下方说明
section_header(slide, '📌 定时任务（交易日09:20自动触发）')
R(slide, MARGIN_L, CONTENT_Y + 1800000, CONTENT_W, 3200000, C_CARD)

cron_items = [
    '① ETF流水线（step0~step8）：全流程约9.5分钟，含申万行业刷新、行情采集、估值合并、低估筛选',
    '② 日报生成推送（09:40）：生成Markdown日报，调用message工具推送微信',
    '③ 数据健康检查（10:00）：检查数据新鲜度，异常时推送微信告警',
    '④ 流水线健康检查（10:05）：检查关键文件+快照连续性',
    '⑤ 微信Session检查（07/11/15/19点）：检查Bot Session有效性，过期时告警',
]
for i, item in enumerate(cron_items):
    TB(slide, MARGIN_L + 20000, CONTENT_Y + 1850000 + i * 600000,
       CONTENT_W - 40000, 550000,
       f'• {item}', Pt(13), color=C_GRAY)

# ============================================================
# PPT Slide 4: 四大数据源
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, '四大数据来源', '覆盖900+只A股ETF · 真实历史分位 + 申万行业估算')

sources = [
    {
        'name': '乐咕乐股',
        'icon': '📈',
        'color': C_GREEN,
        'items': [
            '宽基指数PE/PB历史分位',
            '回溯20年真实历史数据',
            '覆盖：沪深300、中证500、',
            '  创业板、科创50等9个指数',
            '可信度：⭐⭐⭐⭐⭐ 最高',
        ],
        'tag': '✅ 真实分位',
    },
    {
        'name': '申万行业',
        'icon': '🏭',
        'color': C_BLUE,
        'items': [
            '31个申万一级行业官方PE/PB',
            'akshare: sw_index_first_info()',
            '支持穿透估值（ETF持仓×',
            '  行业PE加权计算）',
            '可信度：⭐⭐⭐⭐ 较高',
        ],
        'tag': '⚠️ 估算分位',
    },
    {
        'name': '巨潮资讯',
        'icon': '📋',
        'color': C_SKY,
        'items': [
            'ETF成分股数据',
            '前十大持仓占比',
            '行业分布权重',
            '用于穿透估值计算',
            '可信度：⭐⭐⭐ 参考',
        ],
        'tag': '穿透计算',
    },
    {
        'name': '新浪财经',
        'icon': '💹',
        'color': C_ORANGE,
        'items': [
            '20日日均成交额',
            'akshare: fund_etf_hist_sina()',
            '用于流动性筛选',
            '日均≥1亿（正式池）',
            '可信度：⭐⭐⭐⭐⭐ 准确',
        ],
        'tag': '流动性数据',
    },
]

for i, src in enumerate(sources):
    col = i % 2
    row = i // 2
    l = MARGIN_L + col * (CONTENT_W//2 + GAP//2)
    t = CONTENT_Y + row * 3000000
    w = (CONTENT_W - GAP) // 2
    card(slide, l, t, w, 2700000, f'{src["icon"]} {src["name"]}  {src["tag"]}',
         src['items'], src['color'])

# ============================================================
# PPT Slide 5: 筛选条件
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, '筛选条件详解', '正式池：严格标准 | 关注池：宽松流动性 | 宽基指数可适当放宽')

# 三栏：正式池 / 关注池 / 宽基
pools = [
    {
        'name': '🎯 正式池（严格）',
        'color': C_GREEN,
        'items': [
            '① PE% ≤ 30%（分位数历史低位）',
            '② PB% ≤ 30%（市净率低估）',
            '③ PEG < 1（成长性价比）',
            '④ 日均成交额 ≥ 1亿元（流动性）',
            '',
            '→ PE%或PB%均需满足才进入正式池',
            '→ PEG要求growth_rate>0',
        ],
    },
    {
        'name': '👁️ 关注池（宽松）',
        'color': C_ORANGE,
        'items': [
            '① PE% ≤ 30% 或 PB% ≤ 30%',
            '② 日均成交额 ≥ 300万元',
            '   （正式池的3%，门槛低很多）',
            '③ 不要求 PEG',
            '',
            '→ 主要面向低流动性ETF',
            '→ 提供观察参考，不建议重仓',
        ],
    },
    {
        'name': '📊 宽基指数（特殊）',
        'color': C_BLUE,
        'items': [
            'PE% ≤ 50% 或 PB% ≤ 50%',
            '（放宽条件，因分位基准不同）',
            '',
            '原因：宽基指数真实历史分位',
            '  由乐咕乐股提供，回溯20年',
            '  当前普遍偏高（75%~96%）',
            '  几乎无ETF能满足≤30%阈值',
        ],
    },
]

col_w = (CONTENT_W - 2*GAP) // 3
for i, pool in enumerate(pools):
    l = MARGIN_L + i*(col_w + GAP)
    card(slide, l, CONTENT_Y, col_w, 2400000,
         pool['name'], pool['items'], pool['color'])

# 说明框
R(slide, MARGIN_L, CONTENT_Y + 2500000, CONTENT_W, 2600000, C_CARD)
R(slide, MARGIN_L, CONTENT_Y + 2500000, CONTENT_W, 50000, C_GOLD)
TB(slide, MARGIN_L + 20000, CONTENT_Y + 2550000, CONTENT_W - 40000, 400000,
   '📌 关键指标说明', Pt(14), bold=True, color=C_NAVY)

defs = [
    '• PE%（市盈率分位）：当前PE在历史中的位置。10% = 比历史上90%的时候都便宜',
    '• PB%（市净率分位）：当前PB在历史中的位置。20% = 比历史上80%的时候都便宜',
    '• PEG（市盈率相对盈利增长）：PEG<1 = 盈利增速大于PE，买得划算',
    '• 穿透估值：当ETF无直接PE数据时，用持仓行业×行业PE加权推算（估算数据）',
]
for i, d in enumerate(defs):
    TB(slide, MARGIN_L + 20000, CONTENT_Y + 3050000 + i * 500000,
       CONTENT_W - 40000, 450000, d, Pt(12), color=C_GRAY)

# ============================================================
# PPT Slide 6: 数据质量体系
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, '数据质量体系', '真实分位 vs 估算分位：决定仓位上限，影响组合配置')

# 两栏对比
col_w = (CONTENT_W - GAP) // 2

# 真实数据
R(slide, MARGIN_L, CONTENT_Y, col_w, 3800000, C_CARD)
R(slide, MARGIN_L, CONTENT_Y, col_w, 60000, C_GREEN)
TB(slide, MARGIN_L + 20000, CONTENT_Y + 10000, col_w - 40000, 500000,
   '✅ 真实历史分位（乐咕乐股）', Pt(16), bold=True, color=C_GREEN)
real_items = [
    '数据来源：乐咕乐股（legulegu.com）',
    '历史回溯：20年真实历史PE/PB数据',
    '适用对象：宽基指数ETF（如沪深300ETF）',
    '计算方式：直接查询官方指数PE，',
    '           匹配对应宽基ETF',
    '',
    '优点：数据最可靠，历史规律性强',
    '缺点：覆盖有限，仅9个宽基指数',
    '',
    '当前覆盖：仅1只宽基ETF',
    '  → 红利ETF工银（sz159905）',
    '  → PE%=19.8% PB%=20.2%',
    '',
    '组合上限：balanced模式 ≤40%',
]
for i, item in enumerate(real_items):
    TB(slide, MARGIN_L + 20000, CONTENT_Y + 650000 + i * 260000,
       col_w - 40000, 260000, item, Pt(12), color=C_GRAY)

# 估算数据
R(slide, MARGIN_L + col_w + GAP, CONTENT_Y, col_w, 3800000, C_CARD)
R(slide, MARGIN_L + col_w + GAP, CONTENT_Y, col_w, 60000, C_ORANGE)
TB(slide, MARGIN_L + col_w + GAP + 20000, CONTENT_Y + 10000, col_w - 40000, 500000,
   '⚠️ 估算分位（申万行业/穿透估值）', Pt(16), bold=True, color=C_ORANGE)
est_items = [
    '数据来源：申万行业PE × ETF持仓权重',
    '估算逻辑：ETF申万行业占比 × 行业PE',
    '适用对象：行业/主题ETF（如酒ETF、券商ETF）',
    '',
    '优点：覆盖广，可估算900+只行业ETF',
    '缺点：存在±20%估算误差',
    '       单一行业ETF分位统一（17.7%）',
    '',
    '当前覆盖：正式池36只均为估算数据',
    '  → 酒ETF（sh512690）PE%=19.8%',
    '  → 证券公司ETF  PE%=17.7%（统一值）',
    '',
    '组合上限：balanced模式 ≤30%（已上调）',
]
for i, item in enumerate(est_items):
    TB(slide, MARGIN_L + col_w + GAP + 20000, CONTENT_Y + 650000 + i * 260000,
       col_w - 40000, 260000, item, Pt(12), color=C_GRAY)

# Portfolio Engine 说明
R(slide, MARGIN_L, CONTENT_Y + 3900000, CONTENT_W, 1900000, C_CARD)
R(slide, MARGIN_L, CONTENT_Y + 3900000, CONTENT_W, 50000, C_BLUE)
TB(slide, MARGIN_L + 20000, CONTENT_Y + 3950000, CONTENT_W - 40000, 450000,
   '📊 Portfolio Engine 数据质量仓位控制', Pt(14), bold=True, color=C_BLUE)

pe_items = [
    '• estimated_data_cap 上调（0.15→0.30）：正式池36只行业ETF在组合中最高可占30%总仓位',
    '• real_data_cap（0.40）：真实数据ETF合计上限40%',
    '• 估算数据单只上限：30%÷36只 ≈ 0.83%/只（均分保证分散）',
    '• 风控触发后：若估算合计超30%，自动缩放至上限（当前80%→30%，缩放比例37.5%）',
]
for i, item in enumerate(pe_items):
    TB(slide, MARGIN_L + 20000, CONTENT_Y + 4450000 + i * 330000,
       CONTENT_W - 40000, 330000, item, Pt(12), color=C_GRAY)

# ============================================================
# PPT Slide 7: 为什么只有41只ETF
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, '为什么只有 41 只 ETF？', '41只 = 36只行业ETF（估算） + 1只宽基（真实） + 4只无分位（关注池边缘）')

# 核心问题
R(slide, MARGIN_L, CONTENT_Y, CONTENT_W, 1200000, RGBColor(0xFF, 0xF3, 0xE0))
R(slide, MARGIN_L, CONTENT_Y, 80000, 1200000, C_ORANGE)
TB(slide, MARGIN_L + 120000, CONTENT_Y + 100000, CONTENT_W - 240000, 1000000,
   '💡 核心矛盾：宽基指数ETF当前普遍偏贵，无一满足PE%≤30%的严格标准',
   Pt(16), bold=True, color=C_RED)
TB(slide, MARGIN_L + 120000, CONTENT_Y + 700000, CONTENT_W - 240000, 450000,
   '申万行业PE/PB估算数据可覆盖行业/主题ETF，但估算存在误差，且多只ETF共享同一行业分位值',
   Pt(13), color=C_GRAY, italic=True)

# 三层分析
sections = [
    {
        'title': '❶ 宽基指数ETF → 真实分位（最可靠）',
        'color': C_GREEN,
        'items': [
            '数据来源：乐咕乐股20年历史分位',
            '覆盖范围：仅9个宽基指数（沪深300、中证500、创业板、科创50等）',
            '对应ETF：约169只宽基ETF',
            '现状：当前PE%普遍在 75%~96%（历史高位）',
            '       没有任何一只宽基ETF满足 PE%≤30%',
            '结果：宽基ETF全部被挡在正式池门外',
            '',
            '✅ 唯一例外：红利ETF工银（sz159905）',
            '   PE%=19.8% PB%=20.2% 真实数据',
            '   在关注池（流动性≥300万即可）',
        ],
    },
    {
        'title': '❷ 行业/主题ETF → 申万行业估算分位',
        'color': C_ORANGE,
        'items': [
            '数据来源：申万31个一级行业官方PE × ETF持仓权重',
            '覆盖范围：900+只行业/主题ETF',
            '估算方法：穿透估值（ETF申万行业占比 × 行业PE）',
            '',
            '估算局限：',
            '  ① 同一行业所有ETF分位相同（如7只证券公司ETF均为17.7%）',
            '  ② 估算误差约±20%',
            '  ③ 混合行业ETF分位可能不准确',
            '',
            '正式池结果：36只行业ETF满足 PE%≤30%',
            '  → 酒ETF、证券公司ETF、医疗ETF、红利低波ETF等',
        ],
    },
    {
        'title': '❸ 113只关注池ETF → 流动性不足或其他原因',
        'color': C_BLUE,
        'items': [
            '关注池条件：PE%≤30% 且 日均成交≥300万（正式池门槛的3%）',
            '',
            '不在正式池的原因：',
            '  ① 流动性不足：日均<1亿（正式池门槛）',
            '  ② PEG≥1：增长性价比不足',
            '  ③ 估算数据不稳定：PE%/PB%有缺失值',
            '',
            '典型代表：港股ETF（513090等）、纳斯达克ETF',
            '           日均成交大但申万行业分位不适用',
        ],
    },
]

for i, sec in enumerate(sections):
    t = CONTENT_Y + 1350000 + i * 1750000
    R(slide, MARGIN_L, t, CONTENT_W, 1650000, C_CARD)
    R(slide, MARGIN_L, t, CONTENT_W, 50000, sec['color'])
    TB(slide, MARGIN_L + 20000, t + 80000, CONTENT_W - 40000, 400000,
       sec['title'], Pt(13), bold=True, color=sec['color'])
    for j, item in enumerate(sec['items']):
        TB(slide, MARGIN_L + 40000, t + 480000 + j * 220000,
           CONTENT_W - 80000, 220000, item, Pt(11), color=C_GRAY)

# ============================================================
# PPT Slide 8: 今日低估池正式池（36只）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, '今日低估池：正式池 36 只', f'PE%≤30% · 日均成交≥1亿 · 数据：申万行业估算（⚠️仅供参考）')

# 分类展示
R(slide, MARGIN_L, CONTENT_Y, CONTENT_W, 400000, RGBColor(0xFF, 0xF3, 0xE0))
TB(slide, MARGIN_L + 20000, CONTENT_Y + 80000, CONTENT_W - 40000, 300000,
   '⚠️ 说明：以下36只ETF均为行业/主题ETF，PE%/PB%基于申万行业穿透估算，非真实历史分位',
   Pt(13), color=C_ORANGE, bold=True)

# 分类统计
cats = {
    '🏦 证券公司ETF': [],
    '💊 医疗/生物医药ETF': [],
    '🍷 酒/食品饮料ETF': [],
    '💰 红利/低波ETF': [],
    '🌐 港股/港股通ETF': [],
    '🏨 旅游/其他': [],
}

for etf in formal_pool:
    name = etf.get('name', '')
    code = etf.get('code', '')
    pe = etf.get('pe_percentile', 'N/A')
    pb = etf.get('pb_percentile', 'N/A')
    amt = etf.get('avg_amount_20d', 0) / 1e8
    label = f'{name}({code})  PE%={pe} PB%={pb}  成交{amt:.1f}亿'
    
    if '证券' in name:
        cats['🏦 证券公司ETF'].append(label)
    elif any(k in name for k in ['医疗', '生物', '药', '器械']):
        cats['💊 医疗/生物医药ETF'].append(label)
    elif any(k in name for k in ['酒', '食品']):
        cats['🍷 酒/食品饮料ETF'].append(label)
    elif any(k in name for k in ['红利', '低波', '高股息', '价值']):
        cats['💰 红利/低波ETF'].append(label)
    elif any(k in name for k in ['港股', '香港', 'H股', '恒生']):
        cats['🌐 港股/港股通ETF'].append(label)
    else:
        cats['🏨 旅游/其他'].append(label)

col_w = (CONTENT_W - GAP) // 2
col = 0
row = 0
current_y = CONTENT_Y + 450000

for cat_name, items in cats.items():
    if not items:
        continue
    if col == 2:
        col = 0
        row += 1
    l = MARGIN_L + col * (col_w + GAP)
    t = current_y + row * 2000000
    h = min(len(items) * 280000 + 300000, 1900000)
    
    R(slide, l, t, col_w, h, C_CARD)
    color = C_BLUE if col == 0 else C_ORANGE
    R(slide, l, t, col_w, 50000, color)
    TB(slide, l + 15000, t + 8000, col_w - 30000, 400000,
       cat_name, Pt(12), bold=True, color=C_WHITE)
    for j, item in enumerate(items):
        TB(slide, l + 15000, t + 550000 + j * 280000,
           col_w - 30000, 280000, f'• {item}', Pt(10), color=C_GRAY)
    
    col += 1
    if col == 2:
        col = 0
        row += 1
        current_y = t + h + GAP

# 关注池数据
R(slide, MARGIN_L, CONTENT_Y + 4550000, CONTENT_W, 1200000, C_CARD)
R(slide, MARGIN_L, CONTENT_Y + 4550000, CONTENT_W, 50000, C_SKY)
TB(slide, MARGIN_L + 20000, CONTENT_Y + 4558000, CONTENT_W - 40000, 450000,
   f'📌 关注池（113只）：PE%≤30% · 日均成交≥300万 · 重点推荐：红利ETF工银(sz159905) PE%=19.8% PB%=20.2% ✅真实数据',
   Pt(12), color=C_BLUE)

# ============================================================
# PPT Slide 9: 组合配置
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, 'Portfolio Engine 组合配置建议', '平衡模式 · estimated_data_cap已上调至30%')

# 统计信息
total_w = portfolio_meta.get('total_weight', 0) * 100
cash_w  = portfolio_meta.get('cash_weight', 0) * 100
etf_count = portfolio_meta.get('total_etfs', 0)
warnings = portfolio_meta.get('warnings', [])

R(slide, MARGIN_L, CONTENT_Y, CONTENT_W, 900000, C_CARD)
R(slide, MARGIN_L, CONTENT_Y, CONTENT_W, 50000, C_BLUE)

stats_text = f'总仓位: {total_w:.1f}%   |   现金: {cash_w:.1f}%   |   ETF数量: {etf_count}只   |   警告: {len(warnings)}条'
TB(slide, MARGIN_L + 20000, CONTENT_Y + 80000, CONTENT_W - 40000, 450000,
   stats_text, Pt(16), bold=True, color=C_BLUE)
TB(slide, MARGIN_L + 20000, CONTENT_Y + 520000, CONTENT_W - 40000, 350000,
   f'⚠️ 当前总仓位受限：36只ETF全为估算数据，estimated_data_cap=30%，经风控缩放后总仓位30%',
   Pt(12), color=C_ORANGE, italic=True)

# 正式池 + 关注池分开展示
formal_w = [x for x in portfolio_items if x.get('data_source') != 'legulegu']
real_w   = [x for x in portfolio_items if x.get('data_source') == 'legulegu']

# 如果portfolio里没有数据（全部缩放到30%内），直接用正式池数据
if not portfolio_items:
    # 取正式池前20只按PE%排序
    sorted_formal = sorted(formal_pool, key=lambda x: float(x.get('pe_percentile', 100) or 100))
    display_items = sorted_formal[:20]
else:
    display_items = portfolio_items[:20]

# 两栏：正式池推荐 / 宽基推荐
col_w = (CONTENT_W - GAP) // 2

R(slide, MARGIN_L, CONTENT_Y + 1000000, col_w, 4600000, C_CARD)
R(slide, MARGIN_L, CONTENT_Y + 1000000, col_w, 50000, C_ORANGE)
TB(slide, MARGIN_L + 20000, CONTENT_Y + 1008000, col_w - 40000, 400000,
   f'⚠️ 正式池 TOP15（估算数据 · 申万行业穿透）', Pt(14), bold=True, color=C_ORANGE)

sorted_formal = sorted(formal_pool, key=lambda x: float(x.get('pe_percentile', 100) or 100))
for i, etf in enumerate(sorted_formal[:15], 1):
    pe = etf.get('pe_percentile', 'N/A')
    pb = etf.get('pb_percentile', 'N/A')
    amt = etf.get('avg_amount_20d', 0) / 1e8
    weight = 0.30 / len(sorted_formal[:15]) * 100  # 估算仓位
    TB(slide, MARGIN_L + 20000, CONTENT_Y + 1050000 + i * 280000,
       col_w - 40000, 280000,
       f'{i}. {etf["name"]}({etf["code"]})  PE%={pe} PB%={pb}  成交{amt:.1f}亿  建议~{weight:.1f}%',
       Pt(11), color=C_GRAY)

R(slide, MARGIN_L + col_w + GAP, CONTENT_Y + 1000000, col_w, 4600000, C_CARD)
R(slide, MARGIN_L + col_w + GAP, CONTENT_Y + 1000000, col_w, 50000, C_GREEN)
TB(slide, MARGIN_L + col_w + GAP + 20000, CONTENT_Y + 1008000, col_w - 40000, 400000,
   f'✅ 宽基/真实数据 ETF（推荐重点关注）', Pt(14), bold=True, color=C_GREEN)

# 关注池中带真实数据的
real_etfs = [x for x in watch_list if x.get('pe_percentile') is not None]
for i, etf in enumerate(real_etfs[:10], 1):
    pe = etf.get('pe_percentile', 'N/A')
    pb = etf.get('pb_percentile', 'N/A')
    amt = etf.get('avg_amount_20d', 0) / 1e8
    TB(slide, MARGIN_L + col_w + GAP + 20000, CONTENT_Y + 1050000 + i * 280000,
       col_w - 40000, 280000,
       f'{i}. {etf["name"]}({etf["code"]})  PE%={pe} PB%={pb}  成交{amt:.1f}亿',
       Pt(11), color=C_GRAY)

TB(slide, MARGIN_L + col_w + GAP + 20000, CONTENT_Y + 1050000 + 11 * 280000,
   col_w - 40000, 1000000,
   '📌 注意：宽基ETF当前PE%普遍在75-96%，\n  满足≤30%阈值的极少，建议等待机会',
   Pt(11), color=C_RED, italic=True)

# ============================================================
# PPT Slide 10: 今日日报原样呈现
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, '今日日报（原始输出）', '2026-05-23 · 数据来源：AkShare + 申万行业 + 乐咕乐股')

# 白色内容框
R(slide, MARGIN_L, CONTENT_Y, CONTENT_W, 5450000, C_WHITE)

lines = daily_report.split('\n')
y = CONTENT_Y + 60000
for line in lines:
    if not line.strip():
        y += 150000
        continue
    # 标题行
    if line.startswith('【'):
        color = C_NAVY
        size = Pt(14)
        bold = True
    elif line.startswith('📊') or line.startswith('✅') or line.startswith('⚠️') or line.startswith('📈') or line.startswith('📅') or line.startswith('💡') or line.startswith('🔄'):
        color = C_BLUE
        size = Pt(12)
        bold = True
    elif line.startswith('---'):
        color = C_LGRAY
        size = Pt(10)
        bold = False
    elif line and line[0].isdigit() and '. ' in line[:4]:
        color = C_GRAY
        size = Pt(11)
        bold = False
    else:
        color = C_GRAY
        size = Pt(11)
        bold = False
    
    TB(slide, MARGIN_L + 30000, y, CONTENT_W - 60000, 300000,
       line, size, bold=bold, color=color)
    y += 300000

# ============================================================
# PPT Slide 11: 风险提示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 6858000, C_BG)
title_bar(slide, '风险提示与局限性', '量化辅助工具 · 不构成投资建议 · 数据仅供参考')

warnings_data = [
    {
        'icon': '⚠️',
        'title': '数据局限性',
        'color': C_ORANGE,
        'items': [
            '申万行业估算分位存在±20%误差，不可作为唯一依据',
            '同一行业所有ETF分位相同，无法区分行业内优劣',
            '混合行业ETF穿透估值可能不准确',
            '港股/美股ETF无法用申万行业估算，需其他数据源',
        ],
    },
    {
        'icon': '📉',
        'title': '市场风险',
        'color': C_RED,
        'items': [
            'PE/PB分位低不代表不会继续下跌',
            '低估可能持续数年（如2012-2014年熊市）',
            '行业周期性ETF（证券等）波动极大',
            '建议分批建仓，单只仓位≤20%',
        ],
    },
    {
        'icon': '🔧',
        'title': '模型局限',
        'color': C_BLUE,
        'items': [
            'Portfolio Engine未考虑宏观因素、政策风险',
            '相关性分析基于60天历史，可能不适配极端行情',
            'estimated_data_cap=30%是经验值，非理论最优',
            '建议定期复盘，根据市场变化调整参数',
        ],
    },
    {
        'icon': '✅',
        'title': '正确使用方式',
        'color': C_GREEN,
        'items': [
            '作为选股参考，不作为买卖的唯一依据',
            '重点关注✅真实数据ETF（当前仅红利ETF工银）',
            '行业ETF估算数据仅供辅助参考',
            '严格执行仓位控制，避免过度集中',
        ],
    },
]

col_w = (CONTENT_W - GAP) // 2
for i, w in enumerate(warnings_data):
    col = i % 2
    row = i // 2
    l = MARGIN_L + col * (col_w + GAP)
    t = CONTENT_Y + row * 2800000
    card(slide, l, t, col_w, 2650000, f'{w["icon"]} {w["title"]}',
         w['items'], w['color'])

# 底部声明
R(slide, MARGIN_L, CONTENT_Y + 5700000, CONTENT_W, 600000, RGBColor(0xFF, 0xEE, 0xEE))
TB(slide, MARGIN_L + 30000, CONTENT_Y + 5750000, CONTENT_W - 60000, 500000,
   '⚠️ 免责声明：本系统仅提供量化数据参考，不构成任何投资建议。投资有风险，决策需谨慎，请自行判断并承担后果。',
   Pt(13), color=C_RED, italic=True)

# ============================================================
# 保存
# ============================================================
out_path = os.path.join(OUTPUT_DIR, 'ETF系统架构与数据报告_20260523.pptx')
prs.save(out_path)
print(f'✅ PPT已保存: {out_path}')
