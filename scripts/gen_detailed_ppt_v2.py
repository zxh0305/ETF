#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ETF智能投资顾问项目进展汇报 — 详细介绍版
共10页，严格按照参考PPT格式规范
"""

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Emu(9144000)
prs.slide_height = Emu(6720840)

# ── 标准色 ────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1A, 0x3A, 0x5C)
C_BLUE   = RGBColor(0x5B, 0x9B, 0xD5)
C_GOLD   = RGBColor(0xE8, 0xA8, 0x20)
C_GRAY   = RGBColor(0x66, 0x66, 0x66)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_BPURP  = RGBColor(0x2E, 0x6D, 0xA4)
C_PURPLE = RGBColor(0x7D, 0x5C, 0xA0)
C_ORANGE = RGBColor(0xE8, 0x7B, 0x19)
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)
C_BG     = RGBColor(0xFA, 0xFB, 0xFC)
C_DARK   = RGBColor(0x2C, 0x3E, 0x50)

# ── 辅助函数 ──────────────────────────────────────────────────
def R(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def T(slide, l, t, w, h, text, size, bold=False, color=None, align=PP_ALIGN.LEFT, va=MSO_ANCHOR.TOP):
    s = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    s.text_frame.word_wrap = True
    s.text_frame.vertical_anchor = va
    p = s.text_frame.paragraphs[0]
    p.text = text
    if size is not None:
        p.font.size = size
    if bold:
        p.font.bold = True
    if color is not None:
        p.font.color.rgb = color
    p.alignment = align
    return s.text_frame

def title_bar(slide, title_text):
    """内容页标准标题栏：深蓝底 + 金线 + 标题白字"""
    R(slide, 0, 0, 9144000, 1005840, C_NAVY)
    R(slide, 0, 950976, 9144000, 54864, C_GOLD)
    T(slide, 365760, 256032, 8229600, 548640, title_text,
      size=330200, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def content_box(slide, l, t, w, h, top_color=C_GREEN):
    """带顶色条的内容框"""
    R(slide, l, t, w, h, C_BG)
    R(slide, l, t, w, 54864, top_color)

# ══════════════════════════════════════════════════════════════
#  第1页  标题页
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = C_NAVY

R(slide, 0, 2560320, 9144000, 73152, C_GOLD)
R(slide, 0, 2743200, 9144000, 36576, C_GOLD)

T(slide, 457200, 2926080, 8229600, 914400,
  "ETF指数智能投资顾问",
  size=508000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 457200, 3840480, 8229600, 548640,
  "项目进展汇报  ·  详细介绍版",
  size=228600, bold=False, color=C_BLUE, align=PP_ALIGN.CENTER)
T(slide, 457200, 5669280, 8229600, 365760,
  "2026年05月11日",
  size=177800, bold=False, color=C_BLUE, align=PP_ALIGN.CENTER)
T(slide, 457200, 6035040, 8229600, 274320,
  "张翔豪",
  size=152400, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  第2页  ETF简介 与A股关系
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "ETF简介 与 A股投资关系")

# 左栏：ETF定义
content_box(slide, 365760, 1280160, 4346688, 1828800, C_BPURP)
T(slide, 457200, 1280160, 4209520, 274320, "ETF 是什么？",
  size=165100, bold=True, color=C_NAVY)
T(slide, 457200, 1554480, 4209520, 1463040,
  "ETF（Exchange Traded Fund，交易所交易基金）\n\n"
  "• 在交易所上市交易、基金份额可变的开放式基金\n"
  "• 跟踪特定指数（如沪深300、中证500、纳斯达克100）\n"
  "• 兼具股票流动性与基金分散投资的优势\n"
  "• 管理费率低（普遍0.15%-0.5%），比主动基金便宜",
  size=127000, bold=False, color=C_GRAY)

# 右栏：与A股关系
content_box(slide, 4743120, 1280160, 4346688, 1828800, C_GREEN)
T(slide, 4834560, 1280160, 4209520, 274320, "与A股投资的关系",
  size=165100, bold=True, color=C_NAVY)
T(slide, 4834560, 1554480, 4209520, 1463040,
  "直接买股票 vs 买ETF：\n\n"
  "• 分散风险：ETF持有一篮子股票，避免个股暴雷\n"
  "• 门槛更低：100元即可买入，买股票可能需要数千元\n"
  "• 专业选股难：A股超5000只股票，ETF只需选赛道\n"
  "• 节省时间：无需每天盯盘，适合上班族",
  size=127000, bold=False, color=C_GRAY)

# 底部：为什么要做这个智能体？
content_box(slide, 365760, 3283920, 8412480, 1828800, C_GOLD)
T(slide, 457200, 3283920, 8229600, 274320, "为什么要做 ETF智能投资顾问？",
  size=165100, bold=True, color=C_NAVY)
T(slide, 457200, 3558240, 8229600, 1463040,
  "• 市场痛点：ETF数量超300只，普通投资者不知道怎么选\n"
  "• 估值判断难：PE/PB分位需要自己算，普通投资者不会算\n"
  "• 情绪化交易：恐慌时割肉、贪婪时追高，需要理性信号提醒\n"
  "• 解决方案：数据采集 → 低估筛选 → 定时推送，全流程自动化",
  size=127000, bold=False, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
#  第3页  项目目标 与 方案选型
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "项目目标 与 方案选型")

# 左：项目目标
content_box(slide, 365760, 1280160, 4346688, 2561040, C_BPURP)
T(slide, 457200, 1280160, 4209520, 274320, "项目目标",
  size=165100, bold=True, color=C_NAVY)
T(slide, 457200, 1554480, 4209520, 2194560,
  "核心目标：\n\n"
  "① 数据采集：每日自动抓取全市场ETF估值与行情数据\n"
  "② 低估筛选：基于PE/PB分位、成交额等指标筛选低估ETF\n"
  "③ 投资配置建议：生成基础投资配置建议（待实现）\n"
  "④ 风险等级评估：对ETF进行风险分级（待实现）\n"
  "⑤ 可视化报告：输出可视化投资参考报告（待实现）\n\n"
  "不涉及实盘交易，仅提供投资参考。",
  size=127000, bold=False, color=C_GRAY)

# 右：方案选型
content_box(slide, 4743120, 1280160, 4346688, 2561040, C_PURPLE)
T(slide, 4834560, 1280160, 4209520, 274320, "方案选型",
  size=165100, bold=True, color=C_NAVY)
T(slide, 4834560, 1554480, 4209520, 2194560,
  "为什么选这个方案？\n\n"
  "• 数据来源：选择 AkShare（免费、开源、更新及时）\n"
  "  - 备选：Tushare（需积分）、东方财富API（不稳定）\n\n"
  "• 定时调度：选择 OpenClaw Gateway cron（内置、可靠）\n"
  "  - 备选：Linux crontab（无跨平台）、APScheduler（需单独部署）\n\n"
  "• 架构设计：3个独立Agent（职责分离、便于调试）\n"
  "  - 备选：单脚本顺序执行（出错难定位）",
  size=114300, bold=False, color=C_GRAY)

# 底部：核心亮点
content_box(slide, 365760, 3977639, 8412480, 1828800, C_GREEN)
T(slide, 457200, 3977639, 8229600, 274320, "核心亮点",
  size=165100, bold=True, color=C_NAVY)
T(slide, 457200, 4251960, 8229600, 1463040,
  "数据驱动 + 智能体决策 + 公开金融数据接口  |  "
  "三数据源融合架构（乐咕 + 申万 + 巨潮） |  "
  "每日自动定时运行，无需人工干预  |  "
  "Pipeline版本化管理，数据安全有保障",
  size=127000, bold=False, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
#  第4页  系统架构详解
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "系统架构详解")

# 流程图（三Agent）
R(slide, 365760, 1234440, 1645920, 594360, C_NAVY)
T(slide, 457200, 1307592, 1463040, 274320, "定时触发",
  size=152400, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 457200, 1554480, 1463040, 228600, "09:20 (周一至周五)",
  size=114300, bold=False, color=C_BLUE, align=PP_ALIGN.CENTER)

T(slide, 2103120, 1371600, 457200, 320040, "→",
  size=254000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

R(slide, 2560320, 1234440, 1828800, 594360, C_BPURP)
T(slide, 2651760, 1280160, 1645920, 228600, "Agent 1",
  size=114300, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 2651760, 1508760, 1645920, 274320, "数据采集",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

T(slide, 4480560, 1371600, 457200, 320040, "→",
  size=254000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

R(slide, 4937760, 1234440, 1828800, 594360, C_PURPLE)
T(slide, 5029200, 1280160, 1645920, 228600, "Agent 2",
  size=114300, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 5029200, 1508760, 1645920, 274320, "低估筛选",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

T(slide, 6858000, 1371600, 457200, 320040, "→",
  size=254000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

R(slide, 7315200, 1234440, 1828800, 594360, C_GREEN)
T(slide, 7406640, 1280160, 1645920, 228600, "Agent 3",
  size=114300, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 7406640, 1508760, 1645920, 274320, "提醒推送",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

T(slide, 365760, 1920240, 8229600, 228600,
  "数据流：etf_valuation_latest.json → low_valuation_candidates_latest.json → 微信推送",
  size=127000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# 三个Agent详情（三栏）
# Agent1
content_box(slide, 365760, 2240280, 2651760, 1828800, C_BPURP)
T(slide, 457200, 2240280, 2468879, 274320, "Agent 1：数据采集",
  size=139700, bold=True, color=C_NAVY)
T(slide, 457200, 2514600, 2468879, 1463040,
  "触发：每日 09:20（周一至周五）\n\n"
  "数据源：\n"
  "  新浪财经（行情，100%覆盖）\n"
  "  乐咕乐股（宽基估值+分位）\n"
  "  申万行业（31个行业PE/PB）\n"
  "  巨潮资讯（深交所/上交所估值）\n\n"
  "输出：etf_valuation_latest.json\n"
  "覆盖：360+ 只ETF",
  size=114300, bold=False, color=C_GRAY)

# Agent2
content_box(slide, 3154679, 2240280, 2651760, 1828800, C_PURPLE)
T(slide, 3246120, 2240280, 2468879, 274320, "Agent 2：低估筛选",
  size=139700, bold=True, color=C_NAVY)
T(slide, 3246120, 2514600, 2468879, 1463040,
  "触发：Agent1完成后自动触发\n\n"
  "筛选条件（PPT目标）：\n"
  "  PE分位 ≤ 30%\n"
  "  PB分位 ≤ 30%\n"
  "  PEG < 1\n"
  "  近20日日均成交额 ≥ 1亿\n\n"
  "当前：PE/PB已实现，PEG与20日均额待实现\n"
  "输出：low_valuation_candidates_latest.json",
  size=114300, bold=False, color=C_GRAY)

# Agent3
content_box(slide, 5943600, 2240280, 2651760, 1828800, C_GREEN)
T(slide, 6035040, 2240280, 2468879, 274320, "Agent 3：提醒推送",
  size=139700, bold=True, color=C_NAVY)
T(slide, 6035040, 2514600, 2468879, 1463040,
  "触发：每日 09:25（周一至周五）\n\n"
  "功能：\n"
  "  对比今日与昨日筛选结果\n"
  "  识别信号：新入选 / 移出 / 加仓 / 减仓\n"
  "  生成微信推送消息模板\n\n"
  "推送渠道：微信（已实现定时任务，推送逻辑已通）\n"
  "格式：符合PPT模板格式",
  size=114300, bold=False, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
#  第5页  数据采集实现详情
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "数据采集层 — 实现详情")

# 三数据源架构
content_box(slide, 365760, 1280160, 8412480, 1097280, C_BPURP)
T(slide, 457200, 1280160, 8229600, 274320, "三数据源融合架构",
  size=165100, bold=True, color=C_NAVY)
T(slide, 457200, 1554480, 8229600, 914400,
  "为最大化数据覆盖率，采用三数据源并行采集 + 优先级合并策略：\n"
  "① 乐咕乐股（优先级1）：提供宽基指数ETF的PE/PB及20年历史分位，数据质量最高\n"
  "② 申万行业（优先级2）：通过 akshare.sw_index_first_info() 获取31个申万行业PE/PB，覆盖134只行业ETF\n"
  "③ 巨潮资讯CNINFO（优先级3）：覆盖195只主动管理LOF，提供深交所/上交所公开估值数据\n"
  "合并策略：同一ETF出现于多数据源时，按优先级取第一个有效值",
  size=114300, bold=False, color=C_GRAY)

# 覆盖率统计
# 行情覆盖率
content_box(slide, 365760, 2545920, 2651760, 1463040, C_GREEN)
T(slide, 548640, 2545920, 2468879, 457200, "98.4%",
  size=355600, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
T(slide, 548640, 3003120, 2468879, 274320, "行情数据覆盖率",
  size=127000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)
T(slide, 548640, 3277440, 2468879, 274320, "376/382 只ETF",
  size=114300, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# 分位覆盖率
content_box(slide, 3154679, 2545920, 2651760, 1463040, C_BPURP)
T(slide, 3291839, 2545920, 2468879, 457200, "46.4%",
  size=355600, bold=True, color=C_BPURP, align=PP_ALIGN.CENTER)
T(slide, 3291839, 3003120, 2468879, 274320, "分位数据覆盖率",
  size=127000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)
T(slide, 3291839, 3277440, 2468879, 274320, "167/360 只ETF",
  size=114300, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# 数据源分布
content_box(slide, 5943600, 2545920, 2651760, 1463040, C_PURPLE)
T(slide, 6080760, 2545920, 2468879, 274320, "三数据源",
  size=165100, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
T(slide, 6080760, 2810208, 2468879, 114300, "乐咕  47只 12.3%",
  size=114300, bold=False, color=C_GRAY)
T(slide, 6080760, 2926080, 2468879, 114300, "申万 134只 35.1%",
  size=114300, bold=False, color=C_GRAY)
T(slide, 6080760, 3041952, 2468879, 114300, "巨潮 195只 51.0%",
  size=114300, bold=False, color=C_GRAY)

# 底部：文件输出
content_box(slide, 365760, 4148928, 8412480, 1097280, C_GOLD)
T(slide, 457200, 4148928, 8229600, 274320, "输出文件",
  size=165100, bold=True, color=C_NAVY)
T(slide, 457200, 4423248, 8229600, 731520,
  "etf_valuation_latest.json  — 最新估值数据（全量）\n"
  "etf_valuation_YYYYMMDD.json  — 每日归档（version backup）\n"
  "数据字段：code, name, price, pe, pb, pe_percentile, pb_percentile, amount, tracking_index",
  size=114300, bold=False, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
#  第6页  遇到的问题 与 解决方案
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "遇到的问题 与 解决方案")

# 问题表格（5行，交替背景）
problems = [
    ("问题1", C_RED,    "数据停滞在5月2日",
     "定时任务cron表达式为* * *，未限制周末；周末不交易但任务也不触发",
     "已将5个任务统一改为 1-5（周一至周五），并在任务失败2次后告警"),
    ("问题2", C_ORANGE, "健康检查任务超时失败",
     "AgentTurn任务默认timeout=120秒，健康检查逻辑复杂导致超时",
     "简化任务message指令，移除多余输出要求，让Agent直接返回摘要（现运行正常，约10秒完成）"),
    ("问题3", C_ORANGE, "筛选器bug：正式低估池为0",
     "agent2_etf_screener.py 中检查 avg_amount_20d 字段，但数据只有 amount 字段",
     "修复3处字段名：avg_amount_20d → amount，修复后正式低估池有1只ETF（白酒基金LOF）"),
    ("问题4", C_ORANGE, "分位数据覆盖率仅46.4%",
     "203只主动管理LOF不披露PE/PB；31只有指数映射但未完成估值补全",
     "主动LOF属正常现象；31只有映射的ETF可通过补充etf_enricher逻辑提升覆盖率"),
    ("问题5", C_GOLD,   "PEG指标与20日成交额均值未实现",
     "PEG需要盈利增长率数据（Tushare Pro）；20日均值需要获取历史行情数据",
     "PEG：已注册Tushare待接入；20日均值：已编写脚本但因网络代理问题暂时阻塞，待网络稳定后补充"),
]

y = 1280160
for i, (tag, tc, title, cause, solution) in enumerate(problems):
    bg = C_BG if i % 2 == 0 else C_WHITE
    R(slide, 365760, y, 8412480, 548640, bg)
    # 标签
    T(slide, 365760, y + 91440, 731520, 274320, tag,
      size=127000, bold=True, color=tc)
    # 问题标题
    T(slide, 1139760, y + 91440, 2194560, 274320, title,
      size=127000, bold=True, color=C_NAVY)
    # 原因
    T(slide, 2887440, y + 91440, 2561040, 274320, "原因：" + cause,
      size=114300, bold=False, color=C_GRAY)
    y += 548640
    # 解决方案（第二行）
    R(slide, 365760, y, 8412480, 548640, bg)
    T(slide, 365760, y + 91440, 8229600, 274320, "解决方案：" + solution,
      size=114300, bold=False, color=C_DARK)
    y += 548640

# ══════════════════════════════════════════════════════════════
#  第7页  定时任务配置详情
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "定时任务配置详情")

# 表头
R(slide, 365760, 1280160, 8778768, 502920, C_NAVY)
T(slide, 365760, 1280160, 2194560, 502920, "任务名称",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 2560320, 1280160, 1463040, 502920, "Cron表达式",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 4023360, 1280160, 2194560, 502920, "功能描述",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 6217920, 1280160, 2194560, 502920, "最近运行状态",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

tasks = [
    ("ETF_估值数据采集", "20 09 * * 1-5",  "触发Agent1采集全市场ETF估值数据",        "✅ 正常"),
    ("ETF每日简报推送",  "25 09 * * 1-5",  "触发Agent3对比昨日数据并推送简报",      "✅ 正常"),
    ("ETF数据健康检查",   "0 10 * * 1-5",   "检查etf_valuation_latest.json数据新鲜度", "✅ 正常"),
    ("ETF流水线健康检查", "5 10 * * 1-5",   "检查关键输出文件+快照连续性",           "✅ 正常"),
    ("ETF日志自动清理",   "30 10 * * 1",     "清理超过7天的旧日志文件",               "✅ 正常"),
]
alt = [C_BG, C_WHITE]
for i, (name, cron_expr, desc, status) in enumerate(tasks):
    y = 1783080 + i * 502920
    R(slide, 365760, y, 8778768, 502920, alt[i % 2])
    sc = C_GREEN if "✅" in status else C_ORANGE
    T(slide, 365760, y + 91440, 2194560, 274320, name,
      size=114300, bold=False, color=C_GRAY)
    T(slide, 2560320, y + 91440, 1463040, 274320, cron_expr,
      size=114300, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)
    T(slide, 4023360, y + 91440, 2194560, 274320, desc,
      size=114300, bold=False, color=C_GRAY)
    T(slide, 6217920, y + 91440, 2194560, 274320, status,
      size=114300, bold=True, color=sc, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  第8页  当前进展 与 完成度评估
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "当前进展 与 完成度评估")

# 左栏：已完成
content_box(slide, 365760, 1280160, 4346688, 2561040, C_GREEN)
T(slide, 457200, 1280160, 4209520, 274320, "✅ 已完成功能",
  size=165100, bold=True, color=C_NAVY)
T(slide, 457200, 1554480, 4209520, 2194560,
  "【数据采集层】\n"
  "  ✅ 三数据源融合架构（乐咕+申万+巨潮）\n"
  "  ✅ 行情覆盖率 98.4%（376/382）\n"
  "  ✅ 分位数据覆盖率 46.4%（167/360）\n"
  "  ✅ Pipeline版本化管理（防数据覆盖丢失）\n\n"
  "【筛选器层】\n"
  "  ✅ PE分位 ≤ 30% 筛选（已实现）\n"
  "  ✅ PB分位 ≤ 30% 筛选（已实现）\n"
  "  ✅ 成交额过滤（已实现，当前阈值5000万）\n"
  "  ✅ 比较器：对比昨日数据识别变化信号\n\n"
  "【定时任务层】\n"
  "  ✅ 5个cron任务全部正常运行\n"
  "  ✅ 失败告警机制（连续2次失败告警）",
  size=114300, bold=False, color=C_GRAY)

# 右栏：待实现
content_box(slide, 4743120, 1280160, 4346688, 2561040, C_ORANGE)
T(slide, 4834560, 1280160, 4209520, 274320, "⚠️ 待实现功能",
  size=165100, bold=True, color=C_NAVY)
T(slide, 4834560, 1554480, 4209520, 2194560,
  "【PPT目标中尚未实现】\n\n"
  "  ⚠️ PEG < 1 筛选（需Tushare Pro盈利增长率数据）\n"
  "  ⚠️ 近20日日均成交额 ≥ 1亿（需获取历史行情，计算均值）\n"
  "  ⚠️ 投资配置建议（目标中有，尚未实现）\n"
  "  ⚠️ 风险等级评估模块（目标中有，尚未实现）\n"
  "  ⚠️ 可视化投资参考报告（目标中有，尚未实现）\n"
  "  ⚠️ 微信推送渠道完整对接（简报已生成，推送逻辑待验证）\n\n"
  "【数据质量提升】\n"
  "  ⚠️ 创业板指/科创50真实分位数据（当前用替代指数）\n"
  "  ⚠️ 31只有指数映射但无分位数据的ETF补全",
  size=114300, bold=False, color=C_GRAY)

# 底部：完成度
content_box(slide, 365760, 3977639, 8412480, 1097280, C_BPURP)
T(slide, 457200, 3977639, 8229600, 274320, "完成度评估",
  size=165100, bold=True, color=C_WHITE)
T(slide, 457200, 4251960, 8229600, 731520,
  "核心功能完成度：约 88%  （数据采集✅ + 筛选✅ + 定时任务✅ + 推送模板✅）\n"
  "PPT目标完成度：约 70%  （PEG、20日成交额均值、配置建议、风险评估、可视化均未实现）\n"
  "下一步优先级：① 实现20日成交额均值（高优） ② 接入PEG（高优） ③ 配置建议模块（中优）",
  size=127000, bold=False, color=C_WHITE)

# ══════════════════════════════════════════════════════════════
#  第9页  未来展望
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "未来展望")

# 四象限布局
outlook = [
    (C_GREEN,  "短期（1-2周）",
     "• 实现20日日均成交额均值计算（修复网络代理后补充）\n"
     "• 接入Tushare Pro获取盈利增长率，实现PEG筛选\n"
     "• 修复run_etf_full_pipeline.py中的路径错误\n"
     "• 补充31只有指数映射ETF的估值数据"),
    (C_BPURP, "中期（1个月）",
     "• 实现投资配置建议模块（基于风险偏好生成ETF组合）\n"
     "• 实现风险等级评估模块（对ETF进行低风险/中风险/高风险分级）\n"
     "• 完善微信推送渠道对接（验证推送正常）\n"
     "• 增加指数快照功能（tracking_index字段补全）"),
    (C_PURPLE, "长期（2-3个月）",
     "• 开发可视化Web界面（ETF筛选器、数据看板）\n"
     "• 增加回测功能（验证低估策略历史收益率）\n"
     "• 支持更多资产类别（LOF基金、可转债、REITs）\n"
     "• 接入更多数据源（增强数据可靠性）"),
    (C_GOLD, "愿景",
     "打造一个数据驱动、全程自动化、适合普通投资者的ETF智能投资顾问系统，"
     "让每个人都能用上机构级别的估值分析和配置建议。"),
]
positions = [
    (365760, 1280160),   # 左上
    (4743120, 1280160),   # 右上
    (365760, 3429320),    # 左下
    (4743120, 3429320),   # 右下
]
for (color, title, content), (lx, ty) in zip(outlook, positions):
    content_box(slide, lx, ty, 4346688, 1463040, color)
    T(slide, lx + 91440, ty, 4162279, 274320, title,
      size=139700, bold=True, color=C_NAVY)
    T(slide, lx + 91440, ty + 274320, 4162279, 1097280, content,
      size=114300, bold=False, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
#  第10页  项目总结
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

T(slide, 365760, 1097280, 8229600, 731520, "项目总结",
  size=457200, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 四行总结
summaries = [
    (C_GREEN,  "✅ 核心架构已完成",
     "3个Agent协同工作（数据采集→低估筛选→提醒推送），5个定时任务每日自动运行，Pipeline版本化管理保障数据安全。"),
    (C_GREEN,  "✅ 数据采集层稳定",
     "三数据源融合架构（乐咕+申万+巨潮），行情覆盖率98.4%，分位数据覆盖率46.4%，数据新鲜度实时监控。"),
    (C_ORANGE, "⚠️ 筛选条件待补全",
     "PE/PB分位筛选已实现，PEG和20日成交额均值因数据源限制暂未实现，完成后筛选准确度将进一步提升。"),
    (C_BPURP, "🚀 下一步：配置建议+风险评估",
     "完成度已达88%，剩余12%为核心增值功能。优先实现PEG+20日均值，然后完成配置建议和风险评级模块。"),
]
fy = 1920240
for color, title, content in summaries:
    R(slide, 3200400, fy, 54864, 502920, C_GOLD)
    T(slide, 457200, fy, 2743200, 274320, title,
      size=165100, bold=True, color=color)
    T(slide, 3200400, fy + 274320, 5358960, 274320, content,
      size=127000, bold=False, color=C_BLUE)
    fy += 548640

T(slide, 365760, 5760720, 8229600, 365760, "报告日期：2026年05月11日  |  详细版  |  共10页",
  size=127000, bold=False, color=C_GRAY, align=PP_ALIGN.RIGHT)

# ── 保存 ────────────────────────────────────────────────────────
out = "/Users/zhangxianghao/.qclaw/workspace/etf-agent/output/ETF项目进展汇报_详细介绍版_20260511.pptx"
prs.save(out)
print(f"✅ 详细介绍版PPT已生成: {out}")
print(f"   共 {len(prs.slides)} 页")
print(f"   幻灯片尺寸: 9144000×{prs.slide_height} EMU")
