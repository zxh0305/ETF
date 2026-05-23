#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ETF智能投资顾问项目进展汇报PPT生成器（精美版）
=============================================
包含背景色、形状装饰、详细内容展开
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def set_slide_background(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox_with_format(slide, left, top, width, height, text, font_size=18, bold=False, color=RGBColor(0, 0, 0)):
    """添加格式化的文本框"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    
    return textbox

def add_shape_with_text(slide, shape_type, left, top, width, height, text, fill_color):
    """添加带文字的形状"""
    shape = slide.shapes.add_shape(
        shape_type, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(255, 255, 255)
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return shape

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 定义配色方案
COLOR_PRIMARY = RGBColor(41, 128, 185)      # 蓝色
COLOR_SECONDARY = RGBColor(52, 152, 219)    # 浅蓝
COLOR_ACCENT = RGBColor(231, 76, 60)        # 红色
COLOR_SUCCESS = RGBColor(46, 204, 113)      # 绿色
COLOR_WARNING = RGBColor(241, 196, 15)      # 黄色
COLOR_DARK = RGBColor(44, 62, 80)           # 深蓝灰
COLOR_LIGHT = RGBColor(236, 240, 241)       # 浅灰

# ===== 第1页：精美标题页 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
set_slide_background(slide, COLOR_DARK)

# 添加装饰形状
shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(2))
shape1.fill.solid()
shape1.fill.fore_color.rgb = COLOR_PRIMARY
shape1.line.fill.background()

shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(10), Inches(2))
shape2.fill.solid()
shape2.fill.fore_color.rgb = COLOR_PRIMARY
shape2.line.fill.background()

# 标题
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
tf = title_box.text_frame
tf.vertical_anchor = 1  # 居中

p = tf.paragraphs[0]
p.text = "ETF智能投资顾问项目"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "进展汇报"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# 副标题
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "数据驱动 + 智能体决策\n本地运行 · 免费开源 · 零代码"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

# ===== 第2页：目录（精美版）=====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, COLOR_LIGHT)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "目录"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = COLOR_DARK
p.alignment = PP_ALIGN.LEFT

# 添加装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = COLOR_PRIMARY
line.line.fill.background()

# 目录项（使用形状）
items = [
    ("01", "ETF与A股市场", COLOR_PRIMARY),
    ("02", "项目初衷", COLOR_SECONDARY),
    ("03", "项目目标", COLOR_ACCENT),
    ("04", "已完成功能", COLOR_SUCCESS),
    ("05", "核心Agent实现", COLOR_WARNING),
    ("06", "实际运行效果", COLOR_PRIMARY),
    ("07", "当前差距分析", COLOR_ACCENT),
    ("08", "下一步计划", COLOR_SUCCESS),
    ("09", "总结与展望", COLOR_DARK),
]

y_pos = 1.8
for num, title, color in items:
    # 数字框
    num_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y_pos), Inches(0.6), Inches(0.5))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = color
    num_box.line.fill.background()
    
    tf = num_box.text_frame
    tf.vertical_anchor = 1
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 标题框
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos), Inches(7.5), Inches(0.5))
    tf = title_box.text_frame
    tf.vertical_anchor = 1
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK
    
    y_pos += 0.7

# ===== 第3页：ETF与A股市场（详细版）=====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(255, 255, 255))

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "01 | ETF与A股市场"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

# 添加装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = COLOR_PRIMARY
line.line.fill.background()

# 内容框1：ETF定义
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(4.3), Inches(2.5))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(230, 245, 255)
box1.line.color.rgb = COLOR_SECONDARY

tf = box1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📈 ETF（交易型开放式指数基金）"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "• 在交易所上市交易，可像股票一样买卖"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "• 跟踪特定指数，分散投资风险"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "• 管理费低（通常0.5%以下），透明度高"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "• 2026年A股市场ETF数量：超800只"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK

# 内容框2：ETF与A股关系
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.8), Inches(4.3), Inches(2.5))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(255, 245, 230)
box2.line.color.rgb = COLOR_ACCENT

tf = box2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📊 ETF与A股市场的关系"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "• A股市场规模：超5000只股票，选择困难"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "• ETF覆盖宽基、行业、主题、跨境等"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "• 投资价值：避免个股风险，获取市场平均收益"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "• 适合量化筛选和自动化交易"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK

# 底部总结框
box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(230, 255, 230)
box3.line.color.rgb = COLOR_SUCCESS

tf = box3.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🎯 为何选择ETF作为投资标的？"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_SUCCESS
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "1. 流动性好：日成交亿元级别，进出自由"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "2. 估值透明：PE/PB等估值指标实时可查"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "3. 分散风险：持有一篮子股票，避免个股黑天鹅"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "4. 成本低廉：管理费率通常0.15%-0.5%，远低于主动基金"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_DARK

# ===== 第4页：项目初衷（详细版）=====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(255, 255, 255))

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "02 | 项目初衷"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

# 添加装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = COLOR_PRIMARY
line.line.fill.background()

# 三个框：投资痛点、技术契机、个人需求
boxes_data = [
    ("💡 投资痛点", 
     ["• 人工盯盘耗时耗力，情绪化决策严重", 
      "• ETF数量多（800+只），人工筛选困难", 
      "• 估值数据分散在多个平台，难以快速获取",
      "• 错过最佳买卖时机，收益打折扣"],
     COLOR_ACCENT, Inches(0.5)),
    
    ("🔧 技术契机", 
     ["• OpenClaw（龙虾）智能体开源，自主执行能力强", 
      "• 免费金融数据接口成熟（AkShare/Tushare）", 
      "• 量化投资从机构走向个人投资者",
      "• AI技术普及，降低自动化门槛"],
     COLOR_SUCCESS, Inches(3.5)),
    
    ("🎯 个人需求", 
     ["• 想要一个24小时在线的'赛博员工'", 
      "• 自动监控市场，推送低估机会", 
      "• 数据本地运行，确保隐私安全",
      "• 零成本使用，适合个人投资者"],
     COLOR_WARNING, Inches(6.5)),
]

for title, items, color, x_pos in boxes_data:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.8), Inches(2.8), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = color
    box.line.width = Pt(2)
    
    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    p.space_after = Pt(12)
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(6)

# ===== 第5页：项目目标（详细版）=====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(255, 255, 255))

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "03 | 项目目标"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

# 添加装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = COLOR_PRIMARY
line.line.fill.background()

# 左栏：核心功能目标
left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.2))
left_box.fill.solid()
left_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
left_box.line.color.rgb = COLOR_PRIMARY

tf = left_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "📈 核心功能目标"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY
p.space_after = Pt(15)

p = tf.add_paragraph()
p.text = "1. 数据采集自动化"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "   • 每日09:20自动拉取全市场ETF数据\n   • 覆盖PE/PB/分位/成交额等核心指标\n   • 多数据源融合，确保数据完整性"
p.font.size = Pt(12)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "2. 智能筛选"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "   • 基于PE/PB分位≤30%筛选\n   • 流动性门槛：日均成交≥1亿\n   • PEG<1确保成长价值匹配\n   • 评分系统：多因子加权"
p.font.size = Pt(12)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "3. 自动推送与提醒"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "   • 微信接收每日简报\n   • 变化信号实时提醒（新入选/移出）\n   • 操作建议生成"
p.font.size = Pt(12)
p.font.color.rgb = COLOR_DARK

# 右栏：技术架构目标
right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.8), Inches(4.3), Inches(5.2))
right_box.fill.solid()
right_box.fill.fore_color.rgb = RGBColor(255, 248, 240)
right_box.line.color.rgb = COLOR_ACCENT

tf = right_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "🎯 技术架构目标"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT
p.space_after = Pt(15)

p = tf.add_paragraph()
p.text = "1. 3个独立Agent分工"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "   • Agent1: ETF_估值数据采集\n   • Agent2: ETF_低估筛选\n   • Agent3: ETF_提醒推送"
p.font.size = Pt(12)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "2. 数据驱动决策"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "   • 多数据源融合（乐咕+申万+CNINFO）\n   • 历史分位计算（20年数据）\n   • 实时数据更新（新浪接口）"
p.font.size = Pt(12)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "3. 安全与隐私"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "   • 本地运行，不上传云端\n   • 辅助决策，不自动交易\n   • 数据加密存储"
p.font.size = Pt(12)
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "📈 预期效果"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_DARK
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "   • 不错过任何低估机会\n   • 降低情绪化决策干扰\n   • 提升投资效率和胜率"
p.font.size = Pt(12)
p.font.color.rgb = COLOR_DARK

# 继续添加更多详细页面...
# 由于篇幅限制，我将生成基础框架，然后保存

# ===== 保存文件 =====
output_path = "/Users/zhangxianghao/.qclaw/workspace/etf-agent/output/ETF智能投资顾问项目进展汇报_精美版.pptx"
prs.save(output_path)

print(f"✅ 精美版PPT已生成: {output_path}")
print(f"   当前包含5页精美幻灯片")
print(f"   包含：")
print(f"   - 精美标题页（背景色+装饰）")
print(f"   - 彩色目录页")
print(f"   - ETF与A股市场详解（3个内容框）")
print(f"   - 项目初衷详解（3个框）")
print(f"   - 项目目标详解（双栏布局）")
print(f"\n💡 提示：由于篇幅限制，当前生成了5页示例")
print(f"   如需完整20+页版本，请告知继续生成")
