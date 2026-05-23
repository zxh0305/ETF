#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ETF智能投资顾问项目进展汇报
严格按照参考PPT：精确EMU位置、字体大小、颜色
"""

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Emu(9144000)   # 10"
prs.slide_height = Emu(6720840)   # 7.35"

# ── 精确颜色（与原PPT一致）────────────────────────────────────
C_NAVY   = RGBColor(0x1A, 0x3A, 0x5C)
C_BLUE   = RGBColor(0x5B, 0x9B, 0xD5)
C_GOLD   = RGBColor(0xE8, 0xA8, 0x20)
C_GRAY   = RGBColor(0x66, 0x66, 0x66)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_BPURP  = RGBColor(0x2E, 0x6D, 0xA4)
C_PURPLE = RGBColor(0x7D, 0x5C, 0xA0)
C_ORANGE = RGBColor(0xE8, 0x7B, 0x19)
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)
C_BGBOX  = RGBColor(0xFA, 0xFB, 0xFC)

# ── 辅助函数 ─────────────────────────────────────────────────────
def R(slide, l, t, w, h, color):
    """画填充矩形，无边框"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def T(slide, l, t, w, h, text, size, bold=False, color=None, align=PP_ALIGN.LEFT):
    """添加文本框，返回 text_frame"""
    s = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    s.text_frame.word_wrap = True
    p = s.text_frame.paragraphs[0]
    p.text = text
    if size is not None:
        p.font.size = size if isinstance(size, int) else Pt(size)
    if bold:
        p.font.bold = True
    if color is not None:
        p.font.color.rgb = color
    p.alignment = align
    return s.text_frame

# ══════════════════════════════════════════════════════════════
#  第1页  标题页
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = C_NAVY

# 金色装饰线1  (原 Shape0: left=0, top=2560320, w=9144000, h=73152)
R(slide, 0, 2560320, 9144000, 73152, C_GOLD)
# 金色装饰线2  (原 Shape1: left=0, top=2743200, w=9144000, h=36576)
R(slide, 0, 2743200, 9144000, 36576, C_GOLD)

# 主标题  (原 Shape2: 457200, 2926080, 8229600, 914400)
T(slide, 457200, 2926080, 8229600, 914400,
  "ETF指数智能投资顾问",
  size=508000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 副标题  (原 Shape3: 457200, 3840480, 8229600, 548640)
T(slide, 457200, 3840480, 8229600, 548640,
  "项目进展汇报",
  size=228600, bold=False, color=C_BLUE, align=PP_ALIGN.CENTER)

# 日期  (原 Shape4: 457200, 5669280, 8229600, 365760)
T(slide, 457200, 5669280, 8229600, 365760,
  "2026年05月11日",
  size=177800, bold=False, color=C_BLUE, align=PP_ALIGN.CENTER)

# 作者  (原 Shape5: 457200, 6035040, 8229600, 274320)
T(slide, 457200, 6035040, 8229600, 274320,
  "张翔豪 · ETF指数智能投资顾问智能体",
  size=152400, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  第2页  项目概述
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题栏深蓝  (原 Shape0: 0,0, 9144000,1005840)
R(slide, 0, 0, 9144000, 1005840, C_NAVY)
# 金色细线  (原 Shape1: 0,950976, 9144000,54864)
R(slide, 0, 950976, 9144000, 54864, C_GOLD)
# 标题文字  (原 Shape2: 365760,256032, 8229600,548640)
T(slide, 365760, 256032, 8229600, 548640,
  "项目概述",
  size=330200, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 白色内容背景框  (原 Shape3: 365760,1280160, 8412480,1645920)
R(slide, 365760, 1280160, 8412480, 1645920, C_BGBOX)
# 内容框顶部绿色条  (原 Shape4: 365760,1280160, 8412480,54864)
R(slide, 365760, 1280160, 8412480, 54864, C_GREEN)

# 项目标题  (原 Shape5: 502920,1389888, 8046720,365760)
T(slide, 502920, 1389888, 8046720, 365760,
  "ETF指数智能投资顾问智能体",
  size=228600, bold=True, color=C_NAVY)

# 项目描述  (原 Shape6: 502920,1783080, 8046720,1005840)
T(slide, 502920, 1783080, 8046720, 1005840,
  "搭建ETF投资辅助智能体，抓取主流ETF历史数据、行情分析，生成基础投资配置建议、风险等级评估，输出可视化投资参考报告，不涉及实盘交易。核心亮点是数据驱动+智能体决策，依托公开金融数据接口。",
  size=139700, bold=False, color=C_GRAY)

# ── 三个统计框 ─────────────────────────────────────────────────
# 框1：数据覆盖率  (原 Shape7-11)
R(slide, 365760, 3108960, 2651760, 1280160, C_BGBOX)
R(slide, 365760, 3108960, 2651760, 54864, C_GREEN)
T(slide, 548640, 3246120, 2468879, 457200, "98.4%",
  size=355600, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
T(slide, 548640, 3703320, 2468879, 274320, "数据覆盖率",
  size=139700, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)
T(slide, 548640, 3977639, 2468879, 274320, "376/382只ETF",
  size=127000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# 框2：智能Agent  (原 Shape12-16)
R(slide, 3154679, 3108960, 2651760, 1280160, C_BGBOX)
R(slide, 3154679, 3108960, 2651760, 54864, C_BPURP)
T(slide, 3291839, 3246120, 2468879, 457200, "3",
  size=355600, bold=True, color=C_BPURP, align=PP_ALIGN.CENTER)
T(slide, 3291839, 3703320, 2468879, 274320, "智能Agent",
  size=139700, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)
T(slide, 3291839, 3977639, 2468879, 274320, "采集/筛选/推送",
  size=127000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# 框3：定时任务  (原 Shape17-21)
R(slide, 5943600, 3108960, 2651760, 1280160, C_BGBOX)
R(slide, 5943600, 3108960, 2651760, 54864, C_PURPLE)
T(slide, 6080760, 3246120, 2468879, 457200, "5",
  size=355600, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)
T(slide, 6080760, 3703320, 2468879, 274320, "定时任务",
  size=139700, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)
T(slide, 6080760, 3977639, 2468879, 274320, "每日自动运行",
  size=127000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  第3页  系统架构
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 1005840, C_NAVY)
R(slide, 0, 950976, 9144000, 54864, C_GOLD)
T(slide, 365760, 256032, 8229600, 548640, "系统架构",
  size=330200, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 定时触发框  (原 Shape3-5)
R(slide, 365760, 1234440, 1645920, 594360, C_NAVY)
T(slide, 457200, 1307592, 1463040, 274320, "定时触发",
  size=152400, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 457200, 1554480, 1463040, 228600, "09:20 (1-5)",
  size=114300, bold=False, color=C_BLUE, align=PP_ALIGN.CENTER)

# 箭头1  (原 Shape6)
T(slide, 2103120, 1371600, 457200, 320040, "→",
  size=254000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# Agent1框  (原 Shape7-9)
R(slide, 2560320, 1234440, 1828800, 594360, C_BPURP)
T(slide, 2651760, 1280160, 1645920, 228600, "Agent 1",
  size=114300, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 2651760, 1508760, 1645920, 274320, "ETF_估值数据采集",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 箭头2  (原 Shape10)
T(slide, 4480560, 1371600, 457200, 320040, "→",
  size=254000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# Agent2框  (原 Shape11-13)
R(slide, 4937760, 1234440, 1828800, 594360, C_PURPLE)
T(slide, 5029200, 1280160, 1645920, 228600, "Agent 2",
  size=114300, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 5029200, 1508760, 1645920, 274320, "ETF_低估筛选引擎",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 箭头3  (原 Shape14)
T(slide, 6858000, 1371600, 457200, 320040, "→",
  size=254000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# Agent3框  (原 Shape15-17)
R(slide, 7315200, 1234440, 1828800, 594360, C_GREEN)
T(slide, 7406640, 1280160, 1645920, 228600, "Agent 3",
  size=114300, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 7406640, 1508760, 1645920, 274320, "ETF_提醒推送",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 数据流说明  (原 Shape18)
T(slide, 365760, 1920240, 8229600, 228600,
  "数据流：etf_valuation_latest.json → low_valuation_candidates_latest.json → 微信推送",
  size=127000, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# ── 三个Agent详情框 ─────────────────────────────────────────────
# Agent1详情  (原 Shape19-22)
R(slide, 365760, 2240280, 2651760, 1828800, C_BGBOX)
R(slide, 365760, 2240280, 2651760, 54864, C_BPURP)
T(slide, 457200, 2331720, 2468879, 274320, "Agent 1: 数据采集",
  size=139700, bold=True, color=C_NAVY)
T(slide, 457200, 2606040, 2468879, 1371600,
  "触发时间：每日09:20（周一至周五）\n数据源：AkShare（新浪财经、乐咕乐股、申万、CNINFO）\n输出：etf_valuation_latest.json\n覆盖：全市场360+只ETF实时行情与估值数据",
  size=114300, bold=False, color=C_GRAY)

# Agent2详情  (原 Shape23-26)
R(slide, 3154679, 2240280, 2651760, 1828800, C_BGBOX)
R(slide, 3154679, 2240280, 2651760, 54864, C_PURPLE)
T(slide, 3246120, 2331720, 2468879, 274320, "Agent 2: 筛选引擎",
  size=139700, bold=True, color=C_NAVY)
T(slide, 3246120, 2606040, 2468879, 1371600,
  "触发时间：每日09:20后（依赖Agent1完成）\n筛选条件：PE分位≤30%、PB分位≤30%、成交额≥5000万\n输出：low_valuation_candidates_latest.json\n当前结果：正式低估池1只（白酒基金LOF），观察池359只",
  size=114300, bold=False, color=C_GRAY)

# Agent3详情  (原 Shape27-30)
R(slide, 5943600, 2240280, 2651760, 1828800, C_BGBOX)
R(slide, 5943600, 2240280, 2651760, 54864, C_GREEN)
T(slide, 6035040, 2331720, 2468879, 274320, "Agent 3: 提醒推送",
  size=139700, bold=True, color=C_NAVY)
T(slide, 6035040, 2606040, 2468879, 1371600,
  "触发时间：每日09:25（周一至周五）\n对比昨日数据，识别变化信号\n信号类型：新入选、移出、加仓、减仓\n输出：微信消息模板（符合PPT模板格式）",
  size=114300, bold=False, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
#  第4页  定时任务配置
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 1005840, C_NAVY)
R(slide, 0, 950976, 9144000, 54864, C_GOLD)
T(slide, 365760, 256032, 8229600, 548640, "定时任务配置",
  size=330200, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 表头背景  (原 Shape3: 365760,1280160, 8503919,502920)
R(slide, 365760, 1280160, 8503919, 502920, C_NAVY)

# 表头四列文字  (原 Shape4-7)
# 列位置/宽度：411480/3200400, 3611880/1371600, 4983480/2194560, 7178040/1645920
T(slide, 411480, 1389888, 3200400, 274320, "任务名称",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 3611880, 1389888, 1371600, 274320, "执行时间",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 4983480, 1389888, 2194560, 274320, "功能描述",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, 7178040, 1389888, 1645920, 274320, "状态",
  size=127000, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 表格数据
tasks = [
    ("ETF_估值数据采集",  "09:20 (1-5)", "触发数据采集",             "✅ 正常"),
    ("ETF每日简报推送",  "09:25 (1-5)", "推送简报",                 "✅ 正常"),
    ("ETF数据健康检查",   "10:00 (1-5)", "检查数据新鲜度",           "✅ 正常"),
    ("ETF流水线健康检查", "10:05 (1-5)", "检查关键文件+快照连续性",  "✅ 正常"),
    ("ETF日志自动清理",   "10:30 (每周一)", "清理超过7天的旧日志", "✅ 正常"),
]
alt_colors = [C_BGBOX, C_WHITE]
for i, (name, t, desc, status) in enumerate(tasks):
    y = 1783080 + i * 502920
    bg = alt_colors[i % 2]
    R(slide, 365760, y, 8503919, 502920, bg)
    c = C_GREEN if "✅" in status else C_ORANGE
    T(slide, 411480, y + 91440, 3200400, 274320, name,
      size=114300, bold=False, color=C_GRAY)
    T(slide, 3611880, y + 91440, 1371600, 274320, t,
      size=114300, bold=False, color=C_GRAY)
    T(slide, 4983480, y + 91440, 2194560, 274320, desc,
      size=114300, bold=False, color=C_GRAY)
    T(slide, 7178040, y + 91440, 1645920, 274320, status,
      size=114300, bold=True, color=c)

# ══════════════════════════════════════════════════════════════
#  第5页  实现方案
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 1005840, C_NAVY)
R(slide, 0, 950976, 9144000, 54864, C_GOLD)
T(slide, 365760, 256032, 8229600, 548640, "实现方案",
  size=330200, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 左栏：技术栈  (原 Shape3-6)
R(slide, 365760, 1280160, 3931920, 1645920, C_BGBOX)
R(slide, 365760, 1280160, 3931920, 54864, C_BPURP)
T(slide, 457200, 1371600, 3749039, 274320, "技术栈",
  size=165100, bold=True, color=C_NAVY)
T(slide, 457200, 1645920, 3749039, 1097280,
  "• OpenClaw Gateway (定时调度)\n• Python 3 (脚本处理)\n• AkShare (聚合接口库)\n• 微信消息推送",
  size=127000, bold=False, color=C_GRAY)

# 右栏上部：数据来源  (原 Shape7-10)
R(slide, 365760, 3063240, 3931920, 1463040, C_BGBOX)
R(slide, 365760, 3063240, 3931920, 54864, C_GREEN)
T(slide, 457200, 3154680, 3749039, 274320, "数据来源（通过AkShare封装）",
  size=152400, bold=True, color=C_NAVY)
T(slide, 457200, 3429000, 3749039, 1005840,
  "• 新浪财经 - 行情数据 (100%覆盖)\n• 乐咕乐股 - 宽基指数分位\n• 申万行业 - 行业指数PE/PB\n• 巨潮资讯 - ETF估值数据",
  size=127000, bold=False, color=C_GRAY)

# 右栏：数据流 + 筛选逻辑  (原 Shape11-24)
R(slide, 4434840, 1280160, 3977639, 1280160, C_BGBOX)
R(slide, 4434840, 1280160, 3977639, 54864, C_PURPLE)
T(slide, 4526280, 1371600, 3794759, 274320, "数据流",
  size=165100, bold=True, color=C_NAVY)
T(slide, 4526280, 1645920, 3794759, 320040,
  "AkShare → JSON → 筛选 → 对比 → 推送",
  size=139700, bold=False, color=C_GRAY)
T(slide, 4526280, 1965960, 3794759, 548640,
  "行情：新浪 → 全市场ETF实时价格/成交额\n估值：乐咕/申万/巨潮 → PE/PB/分位数据",
  size=114300, bold=False, color=C_GRAY)

# 筛选逻辑  (原 Shape17-24)
R(slide, 4434840, 2697479, 3977639, 1828800, C_BGBOX)
R(slide, 4434840, 2697479, 3977639, 54864, C_GOLD)
T(slide, 4526280, 2788920, 3794759, 274320, "筛选逻辑",
  size=165100, bold=True, color=C_NAVY)

fy = 3108960
filters = [
    ("PE/PB分位 ≤ 30%",  "当前估值低于历史70%的时间，估值处于低位"),
    ("日均成交额 ≥ 1亿",  "流动性充足，便于买卖（当前使用当日成交额≥5000万，待优化为20日均值）"),
    ("PEG < 1",           "盈利增长率高于市盈率，成长性被低估（注：当前未实现，待接入Tushare Pro）"),
]
for cond, desc in filters:
    T(slide, 4526280, fy, 3794759, 228600, cond,
      size=127000, bold=True, color=C_BPURP)
    fy += 228600
    T(slide, 4526280, fy, 3794759, 228600, f"→ {desc}",
      size=114300, bold=False, color=C_GRAY)
    fy += 304800

# ══════════════════════════════════════════════════════════════
#  第6页  当前进展
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
R(slide, 0, 0, 9144000, 1005840, C_NAVY)
R(slide, 0, 950976, 9144000, 54864, C_GOLD)
T(slide, 365760, 256032, 8229600, 548640, "当前进展",
  size=330200, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 左：数据覆盖率  (原 Shape3-9)
R(slide, 365760, 1280160, 3931920, 1554480, C_BGBOX)
R(slide, 365760, 1280160, 3931920, 54864, C_GREEN)
T(slide, 457200, 1389888, 3749039, 274320, "数据覆盖率",
  size=165100, bold=True, color=C_NAVY)
T(slide, 457200, 1691639, 1828800, 457200, "98.4%",
  size=355600, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
T(slide, 2011680, 1783080, 2103120, 274320, "376/382 只ETF（行情数据）",
  size=139700, bold=False, color=C_GRAY)
T(slide, 457200, 2194560, 3749039, 228600, "注：6只纯主动基金无法覆盖",
  size=114300, bold=False, color=C_GRAY)
# 分位覆盖率补充
T(slide, 457200, 2423160, 1828800, 457200, "46.4%",
  size=355600, bold=True, color=C_BPURP, align=PP_ALIGN.CENTER)
T(slide, 2011680, 2514600, 2103120, 274320, "167/360 只（分位数据）",
  size=139700, bold=False, color=C_GRAY)
T(slide, 457200, 2844659, 3749039, 228600, "注：203只主动管理LOF无PE/PB，31只有指数映射可补充",
  size=114300, bold=False, color=C_GRAY)

# 右：数据源分布  (原 Shape10-13)
R(slide, 4434840, 1280160, 3931920, 1554480, C_BGBOX)
R(slide, 4434840, 1280160, 3931920, 54864, C_BPURP)
T(slide, 4526280, 1389888, 3749039, 274320, "数据源分布（三数据源融合架构）",
  size=165100, bold=True, color=C_NAVY)
T(slide, 4526280, 1691639, 3749039, 914400,
  "申万行业: 134只 (35.1%)\n巨潮资讯: 195只 (51.0%)\n乐咕乐股: 47只 (12.3%)",
  size=139700, bold=False, color=C_GRAY)

# 底部：当前问题与解决方案  (原 Shape14-26)
R(slide, 365760, 2971799, 8412480, 2926080, C_BGBOX)
R(slide, 365760, 2971799, 8412480, 54864, C_RED)
T(slide, 457200, 3063239, 8229600, 274320, "当前问题与解决方案",
  size=165100, bold=True, color=C_RED)

fy = 3428999
issues = [
    ("✅ 已解决", C_GREEN,  "数据停滞问题",  "定时任务已更新为周一至周五（1-5），5月3日后数据已正常采集"),
    ("✅ 已解决", C_GREEN,  "健康检查超时",  "简化任务message指令，避免超时（健康检查任务现正常运行）"),
    ("✅ 已解决", C_GREEN,  "筛选器bug",    "修复avg_amount_20d字段名错误，正式低估池现正常显示1只ETF"),
    ("⚠️ 待实现", C_ORANGE, "20日成交额均值", "需获取近20日历史数据，计算均值后调整阈值为1亿元"),
    ("⚠️ 待实现", C_ORANGE, "PEG指标",       "需注册Tushare Pro获取盈利增长率数据，计算PEG=PE/增长率"),
]
for status, sc, issue, solution in issues:
    T(slide, 457200, fy, 1371600, 274320, status,
      size=127000, bold=True, color=sc)
    T(slide, 1828800, fy, 1371600, 274320, issue,
      size=127000, bold=True, color=C_NAVY)
    T(slide, 3200400, fy, 5003280, 274320, solution,
      size=114300, bold=False, color=C_GRAY)
    fy += 457200

# ══════════════════════════════════════════════════════════════
#  第7页  项目总结
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题（无蓝底，直接白字大号，背景透明）
T(slide, 365760, 1097280, 8229600, 731520, "项目总结",
  size=457200, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 四个要点（纵向排列，左侧橙色竖线装饰）
summaries = [
    (C_GREEN,  "✅ 核心功能",  "3个Agent协同工作（数据采集、低估筛选、提醒推送），5个定时任务正常运行，微信推送正常，按PPT模板格式输出"),
    (C_GREEN,  "✅ 数据覆盖",  "行情数据覆盖率98.4%（376/382），分位数据46.4%（167/360），三数据源架构稳定运行"),
    (C_GREEN,  "✅ 完整流程",  "每日自动采集、筛选、推送全流程打通，Pipeline版本化管理，数据新鲜度实时监控"),
    (C_ORANGE, "⚠️ 待解决",   "完成度88%。剩余12%：20日成交额均值（高优）、PEG指标（高优）、风险评估（中优）、配置建议（中优）"),
]
fy = 1920240
for color, title, content in summaries:
    # 装饰竖线
    R(slide, 3200400, fy, 54864, 502920, C_GOLD)
    T(slide, 457200, fy, 2743200, 274320, title,
      size=165100, bold=True, color=color)
    T(slide, 3200400, fy + 274320, 5358960, 274320, content,
      size=127000, bold=False, color=C_BLUE)
    fy += 548640

# 日期  (原 Shape10: 365760,5760720, 8229600,365760)
T(slide, 365760, 5760720, 8229600, 365760, "报告日期：2026年05月11日",
  size=152400, bold=False, color=C_GRAY, align=PP_ALIGN.RIGHT)

# ── 保存 ────────────────────────────────────────────────────────
out = "/Users/zhangxianghao/.qclaw/workspace/etf-agent/output/ETF项目进展汇报_20260511.pptx"
prs.save(out)
print(f"✅ PPT已生成: {out}")
print(f"   共 {len(prs.slides)} 页")
print(f"   幻灯片尺寸: {prs.slide_width}×{prs.slide_height} EMU")
