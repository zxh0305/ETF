#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ETF智能投资顾问项目进展汇报PPT生成器
=======================================
按照"龙虾驱动的ETF投资智能体.pptx"格式
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for item in content:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """添加两栏内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    # 左栏
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = left.text_frame
    tf.word_wrap = True
    for item in left_content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    # 右栏
    right = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    tf = right.text_frame
    tf.word_wrap = True
    for item in right_content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    return slide

# ===== 开始创建PPT =====

# 第1页：标题页
add_title_slide(prs, 
    "ETF智能投资顾问项目", 
    "进展汇报\n\n数据驱动 + 智能体决策\n本地运行 · 免费开源 · 零代码"
)

# 第2页：目录
add_content_slide(prs, "目录", [
    "01  ETF与A股市场",
    "02  项目初衷",
    "03  项目目标",
    "04  已完成功能",
    "05  核心Agent实现",
    "06  实际运行效果",
    "07  当前差距分析",
    "08  下一步计划",
    "09  总结与展望"
])

# 第3页：ETF与A股市场
add_content_slide(prs, "ETF与A股市场", [
    "📊 ETF（交易型开放式指数基金）",
    "   - 在交易所上市交易，可像股票一样买卖",
    "   - 跟踪特定指数，分散投资风险",
    "   - 管理费低，透明度高",
    "",
    "📈 ETF与A股市场的关系",
    "   - A股市场规模：超5000只股票，选择困难",
    "   - ETF数量：超800只（2026年），覆盖宽基、行业、主题",
    "   - 投资价值：避免个股风险，获取市场平均收益",
    "",
    "🎯 为何选择ETF？",
    "   - 流动性好（日成交亿元级别）",
    "   - 估值透明（PE/PB实时可查）",
    "   - 适合量化筛选和自动化交易"
])

# 第4页：项目初衷
add_content_slide(prs, "项目初衷", [
    "💡 为什么要做这个项目？",
    "",
    "01  投资痛点",
    "   - 人工盯盘耗时耗力，情绪化决策",
    "   - ETF数量多，筛选困难",
    "   - 估值数据分散，难以快速获取",
    "",
    "02  技术契机",
    "   - OpenClaw（龙虾）智能体开源，自主执行能力强",
    "   - 免费金融数据接口成熟（AkShare/Tushare）",
    "   - 量化投资从机构走向个人",
    "",
    "03  个人需求",
    "   - 想要一个24小时在线的'赛博员工'",
    "   - 自动监控市场，推送低估机会",
    "   - 数据本地运行，确保隐私安全"
])

# 第5页：项目目标
add_two_column_slide(prs, "项目目标", 
    [
        "📊 核心功能目标",
        "1. 数据采集自动化",
        "   - 每日09:20自动拉取全市场ETF数据",
        "   - 覆盖PE/PB/分位/成交额等核心指标",
        "",
        "2. 智能筛选",
        "   - 基于PE/PB分位≤30%筛选",
        "   - 流动性门槛：日均成交≥1亿",
        "   - PEG<1确保成长价值匹配",
        "",
        "3. 自动推送",
        "   - 微信接收每日简报",
        "   - 变化信号实时提醒"
    ],
    [
        "🎯 技术架构目标",
        "1. 3个独立Agent分工",
        "   - Agent1: 数据采集",
        "   - Agent2: 智能筛选",
        "   - Agent3: 提醒推送",
        "",
        "2. 数据驱动决策",
        "   - 多数据源融合（乐咕+申万+CNINFO）",
        "   - 历史分位计算（20年数据）",
        "",
        "3. 安全与隐私",
        "   - 本地运行，不上传云端",
        "   - 辅助决策，不自动交易",
        "",
        "📈 预期效果",
        "   - 不错过任何低估机会",
        "   - 降低情绪化决策干扰",
        "   - 提升投资效率和胜率"
    ]
)

# 第6页：已完成功能
add_content_slide(prs, "已完成功能（完成度88%）", [
    "✅ 01  数据采集层（100%）",
    "   - AkShare接口接入，每日09:20自动采集",
    "   - 三数据源架构：乐咕乐股+申万行业+CNINFO",
    "   - 数据覆盖率：46.4%（167/360只ETF有分位数据）",
    "",
    "✅ 02  智能筛选层（70%）",
    "   - PE/PB分位≤30%筛选规则已实现",
    "   - 流动性筛选：当日成交额≥5000万",
    "   - 评分系统：PE分位40%+PB分位30%+流动性30%",
    "",
    "✅ 03  输出推送层（100%）",
    "   - 每日09:25自动推送微信",
    "   - 推送格式已按PPT模板优化",
    "   - 变化信号识别（新入选/移出）"
])

# 第7页：核心Agent实现
add_two_column_slide(prs, "核心Agent实现",
    [
        "🤖 Agent1: ETF_估值数据采集",
        "   - 定时任务：每日09:20（周一至周五）",
        "   - 数据输出：etf_valuation_latest.json",
        "   - 健康监控：10:00检查数据新鲜度",
        "",
        "🤖 Agent2: ETF_低估筛选",
        "   - 输入：估值数据（360只ETF）",
        "   - 输出：low_valuation_candidates_latest.json",
        "   - 正式低估池：1只（白酒基金LOF）",
        "   - 观察池：359只",
        "",
        "🤖 Agent3: ETF_提醒推送",
        "   - 定时任务：每日09:25",
        "   - 推送渠道：微信（openclaw-weixin）",
        "   - 推送内容：今日符合条件+重点关注+变化提醒"
    ],
    [
        "📊 数据流转示意图",
        "   09:20 数据采集",
        "   ↓",
        "   估值补全（三数据源）",
        "   ↓",
        "   09:25 低估筛选",
        "   ↓",
        "   生成简报+变化报告",
        "   ↓",
        "   微信推送",
        "",
        "🔧 定时任务清单",
        "   - ETF数据采集（09:20，1-5）",
        "   - ETF简报推送（09:25，1-5）",
        "   - ETF数据健康检查（10:00，1-5）",
        "   - ETF流水线健康检查（10:05，1-5）",
        "   - ETF日志自动清理（10:30，每周一）"
    ]
)

# 第8页：实际运行效果
add_content_slide(prs, "实际运行效果", [
    "📱 微信推送示例（2026-05-11）",
    "",
    "【ETF低估智能体 | 每日09:20播报】",
    "日期：2026-05-11",
    "",
    "今日符合条件共：1只",
    "",
    "重点关注：",
    "1. 白酒基金LOF（sz161725） PE分位21.5% PB分位21.7% 成交0.68亿",
    "",
    "变化提醒：",
    "🆕 白酒基金LOF新入选低估池",
    "📤 嘉实原油LOF移出观察池",
    "",
    "操作建议：",
    "白酒基金LOF估值偏低，建议等待确认后关注",
    "",
    "—",
    "⚠️ 风险提示：仅供参考，不构成投资建议"
])

# 第9页：当前差距分析
add_two_column_slide(prs, "当前差距分析（完成度88%）",
    [
        "❌ 主要差距1：PEG指标未实现",
        "   - PPT要求：PEG<1（成长价值匹配）",
        "   - 当前状态：完全没有PEG数据",
        "   - 影响：可能选中价值陷阱",
        "   - 解决方案：注册Tushare Pro获取盈利增长率",
        "",
        "❌ 主要差距2：20日成交额均值未实现",
        "   - PPT要求：近20日日均≥1亿",
        "   - 当前状态：使用当日成交额，阈值5000万",
        "   - 影响：筛选结果不稳定",
        "   - 解决方案：计算20日均值，调整阈值为1亿"
    ],
    [
        "✅ 已完成功能",
        "   - 数据采集自动化（09:20定时）",
        "   - PE/PB分位筛选（≤30%）",
        "   - 微信推送（09:25定时）",
        "   - 推送模板（符合PPT格式）",
        "   - 定时任务+健康检查+告警",
        "",
        "📊 完成度评估",
        "   - Agent1: 数据采集 100%",
        "   - Agent2: 低估筛选 70%（缺PEG）",
        "   - Agent3: 提醒推送 100%",
        "   - 推送格式 100%",
        "",
        "🎯 总体完成度：88%"
    ]
)

# 第10页：下一步计划
add_content_slide(prs, "下一步计划", [
    "🚀 优先级排序",
    "",
    "🔴 高优先级（Week 1）",
    "   1. 实现PEG指标 - 工作量2-3天",
    "      - 注册Tushare Pro，获取盈利增长率数据",
    "      - 计算PEG = PE / 盈利增长率",
    "      - 增加筛选条件：PEG<1",
    "",
    "   2. 实现20日成交额均值 - 工作量1天",
    "      - 获取近20日成交额数据",
    "      - 计算均值，调整阈值为1亿",
    "",
    "🟡 中优先级（Week 2）",
    "   3. 风险评估模块 - 工作量1-2天",
    "   4. 配置建议模块 - 工作量2-3天",
    "",
    "🟢 低优先级（Week 3+）",
    "   5. 可视化报告生成器",
    "   6. Web界面部署（可选）"
])

# 第11页：总结与展望
add_two_column_slide(prs, "总结与展望",
    [
        "📊 项目总结",
        "   - 核心功能已基本完成（88%）",
        "   - 数据采集→筛选→推送全流程已打通",
        "   - 定时任务稳定运行，微信推送正常",
        "",
        "💡 核心亮点",
        "   - 数据驱动：多数据源融合，覆盖46.4%",
        "   - 智能体决策：OpenClaw自主执行",
        "   - 本地安全运行：隐私绝对保障",
        "   - 免费开源：零成本使用",
        "",
        "⚠️ 当前限制",
        "   - PEG指标缺失（需Tushare Pro）",
        "   - 20日成交额均值未实现",
        "   - 部分ETF无PE/PB数据（主动LOF）"
    ],
    [
        "🚀 未来展望",
        "   1. AI智能预测",
        "      - 引入机器学习模型，预测ETF未来走势",
        "      - 提升筛选胜率",
        "",
        "   2. 多因子筛选模型",
        "      - 构建更复杂的多因子模型",
        "      - 多维度评估资产，提升准确性",
        "",
        "   3. 策略分享社区",
        "      - 建立开放社区，允许用户分享策略",
        "      - 实现智慧共享",
        "",
        "   4. 实盘对接（可选）",
        "      - 模拟盘测试",
        "      - 用户自主选择是否实盘"
    ]
)

# 第12页：结束页
add_title_slide(prs, 
    "感谢聆听", 
    "ETF智能投资顾问项目进展汇报\n\n本地安全 · 免费 · 零代码\n\n祝你投资顺利，风控为先"
)

# 保存文件
output_path = "/Users/zhangxianghao/.qclaw/workspace/etf-agent/output/ETF智能投资顾问项目进展汇报.pptx"
prs.save(output_path)

print(f"✅ PPT已生成: {output_path}")
print(f"   共12页，包含：")
print(f"   - ETF与A股市场介绍")
print(f"   - 项目初衷与目标")
print(f"   - 已完成功能详解")
print(f"   - 实际运行效果展示")
print(f"   - 差距分析与下一步计划")
