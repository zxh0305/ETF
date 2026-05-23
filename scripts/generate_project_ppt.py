#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ETF指数智能投资顾问项目汇报PPT生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """创建项目汇报PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 定义明亮配色（用户偏好）
    COLOR_PRIMARY = RGBColor(41, 128, 185)      # 蓝色
    COLOR_SECONDARY = RGBColor(52, 152, 219)    # 浅蓝
    COLOR_ACCENT = RGBColor(231, 76, 60)         # 红色
    COLOR_SUCCESS = RGBColor(46, 204, 113)       # 绿色
    COLOR_WARNING = RGBColor(241, 196, 15)       # 黄色
    COLOR_DARK = RGBColor(44, 62, 80)            # 深蓝灰
    
    # ==================== 第1页：封面 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "ETF指数智能投资顾问项目"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "数据驱动 + 智能体决策 · 全流程自动化"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "2026年5月22日"
    p = date_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(127, 140, 141)
    p.alignment = PP_ALIGN.CENTER
    
    # ==================== 第2页：项目概述 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "项目概述"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 三Agent架构
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    text_content = [
        ("三Agent架构", Pt(20), True, COLOR_DARK),
        ("• 数据采集Agent：每日09:00自动触发，采集全市场ETF行情数据", Pt(16), False, COLOR_DARK),
        ("• 筛选Agent：基于PE/PB分位、PEG、流动性筛选低估ETF", Pt(16), False, COLOR_DARK),
        ("• 提醒Agent：对比昨日与今日数据，识别交易信号", Pt(16), False, COLOR_DARK),
        ("", Pt(16), False, COLOR_DARK),
        ("核心亮点", Pt(20), True, COLOR_DARK),
        ("• 数据驱动：四数据源架构，覆盖360只ETF", Pt(16), False, COLOR_DARK),
        ("• 智能体决策：自动化全流程，无需人工干预", Pt(16), False, COLOR_DARK),
        ("• 真实数据优先：仅使用真实历史分位给出投资建议", Pt(16), False, COLOR_DARK),
        ("• 数据透明化：标注数据来源与可信度（✅真实/⚠️估算/❓未知）", Pt(16), False, COLOR_DARK),
    ]
    
    for text, size, bold, color in text_content:
        p = content_frame.add_paragraph() if content_frame.paragraphs else content_frame.paragraphs[0]
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6)
    
    # ==================== 第3页：数据获取（四大数据源） ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "数据获取 · 四大数据源"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 数据源表格
    sources = [
        ("数据源", "获取内容", "覆盖范围", "数据年限"),
        ("乐咕乐股", "宽基指数PE/PB历史分位", "12个宽基指数", "11-21年"),
        ("申万行业", "31个一级行业PE/PB/股息率", "行业/主题ETF估算", "实时"),
        ("巨潮资讯", "国证行业指数估值", "行业ETF补充", "实时"),
        ("新浪财经", "ETF实时行情+20日均成交额", "全市场ETF", "实时"),
    ]
    
    # 创建表格
    rows, cols = len(sources), 4
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 填充表格
    for i, row_data in enumerate(sources):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            # 设置格式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14 if i == 0 else 12)
                paragraph.font.bold = (i == 0)
                paragraph.font.color.rgb = RGBColor(255, 255, 255) if i == 0 else COLOR_DARK
                paragraph.alignment = PP_ALIGN.CENTER
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_PRIMARY
    
    # ==================== 第4页：数据处理流程 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "数据处理流程 · Step0 → Step8"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 流程步骤
    steps = [
        ("Step0", "申万行业估值刷新", "sw_industry_valuation.py", "非必需"),
        ("Step1", "行情数据采集", "etf_data_collector.py", "必需"),
        ("Step2", "估值补全", "etf_valuation_enricher.py", "必需"),
        ("Step2a", "穿透估值计算+合并", "penetration_valuation_v1_real.py", "非必需"),
        ("Step3", "低估筛选", "agent2_etf_screener.py", "必需"),
        ("Step4", "结果对比", "agent3_etf_comparator.py", "非必需"),
        ("Step5", "微信日报生成", "agent3_daily_report.py", "必需"),
        ("Step6-8", "归档+分位计算+趋势分析", "archive_index_snapshot.py等", "非必需"),
    ]
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for step_id, step_name, script, required in steps:
        p = content_frame.add_paragraph()
        p.text = f"【{step_id}】{step_name} ({'✅必需' if required == '必需' else '⚠️非必需'})"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY if required == '必需' else COLOR_WARNING
        p.space_after = Pt(3)
        
        p = content_frame.add_paragraph()
        p.text = f"  脚本: {script}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(6)
    
    # ==================== 第5页：为什么做数据穿透 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "为什么做数据穿透估值？"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    text_content = [
        ("问题：行业/主题ETF缺乏真实历史分位", Pt(18), True, COLOR_ACCENT),
        ("• 宽基ETF：有真实历史分位（乐咕乐股，12个指数）", Pt(14), False, COLOR_DARK),
        ("• 行业/主题ETF：无自身历史数据，只能用申万行业估算", Pt(14), False, COLOR_DARK),
        ("• 当前覆盖率：仅2.8%（1/36）的正式池ETF有真实数据", Pt(14), False, COLOR_ACCENT),
        ("", Pt(14), False, COLOR_DARK),
        ("解决方案：穿透估值", Pt(18), True, COLOR_SUCCESS),
        ("• 核心逻辑：ETF持仓 → 成分股权重 → 个股PE/PB → 加权计算", Pt(14), False, COLOR_DARK),
        ("• 实现方式：基于ETF季报持仓 × 申万行业PE/PB估算", Pt(14), False, COLOR_DARK),
        ("• 当前效果：处理897只ETF，成功269只（成功率30%）", Pt(14), False, COLOR_DARK),
        ("• 局限性：非ETF自身历史分位，仅为持仓加权估算", Pt(14), False, COLOR_WARNING),
    ]
    
    for text, size, bold, color in text_content:
        p = content_frame.add_paragraph() if text else content_frame.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6)
    
    # ==================== 第6页：当前数据覆盖率 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "当前数据覆盖率"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 覆盖率数据
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    text_content = [
        ("真实数据（✅）", Pt(20), True, COLOR_SUCCESS),
        ("• 宽基指数：12个（沪深300、中证500等）", Pt(16), False, COLOR_DARK),
        ("• 数据来源：乐咕乐股（第三方，但为真实历史分位）", Pt(16), False, COLOR_DARK),
        ("• 数据年限：11-21年历史", Pt(16), False, COLOR_DARK),
        ("• 正式池覆盖：仅1只（sz159905 深证红利）", Pt(16), False, COLOR_ACCENT),
        ("", Pt(16), False, COLOR_DARK),
        ("估算数据（⚠️）", Pt(20), True, COLOR_WARNING),
        ("• 申万行业估算：31个一级行业Top20成分股加权", Pt(16), False, COLOR_DARK),
        ("• 穿透估值：269只ETF（成功率30%）", Pt(16), False, COLOR_DARK),
        ("• 正式池覆盖：35/36只（约97%）", Pt(16), False, COLOR_ACCENT),
        ("", Pt(16), False, COLOR_DARK),
        ("目标：提升至80%+真实数据覆盖率", Pt(18), True, COLOR_PRIMARY),
    ]
    
    for text, size, bold, color in text_content:
        p = content_frame.add_paragraph() if text else content_frame.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6)
    
    # ==================== 第7页：定时任务与推送 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "定时任务与推送时间线"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 时间线
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    text_content = [
        ("01:30 凌晨", Pt(16), True, COLOR_DARK),
        ("• ETF穿透估值计算（主任务）", Pt(14), False, COLOR_DARK),
        ("", Pt(14), False, COLOR_DARK),
        ("08:00 早上", Pt(16), True, COLOR_DARK),
        ("• ETF穿透估值计算（备用任务，双重保险）", Pt(14), False, COLOR_DARK),
        ("", Pt(14), False, COLOR_DARK),
        ("09:00 早上", Pt(16), True, COLOR_SUCCESS),
        ("• ETF全流程启动（数据采集→筛选→日报生成）", Pt(14), False, COLOR_DARK),
        ("• 预计09:21完成", Pt(14), False, COLOR_DARK),
        ("", Pt(14), False, COLOR_DARK),
        ("09:25 早上", Pt(16), True, COLOR_SUCCESS),
        ("• ETF每日简报推送（webchat窗口）", Pt(14), False, COLOR_DARK),
        ("• 🎉 09:30前完成推送！", Pt(14), False, COLOR_SUCCESS),
        ("", Pt(14), False, COLOR_DARK),
        ("其他监控任务", Pt(16), True, COLOR_DARK),
        ("• 10:00 ETF数据健康检查 | 10:05 ETF流水线健康检查", Pt(14), False, COLOR_DARK),
        ("• 10:30 每周一日志自动清理 | 07/11/15/19点微信Session健康检查", Pt(14), False, COLOR_DARK),
    ]
    
    for text, size, bold, color in text_content:
        p = content_frame.add_paragraph() if text else content_frame.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6)
    
    # ==================== 第8页：推送格式详解 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "推送格式详解（微信日报）"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 推送格式示例
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # 示例格式（简化版）
    example_text = """📊 ETF每日简报 2026-05-22

【市场温度】温和低估
【投资建议】
✅ 可买清单（数据真实）:
  1. sz159905 深证红利 ✅ 极度低估 PE%=5.2 PB%=8.1 成交额=3.2亿
  2. sh512880 证券ETF ⚠️ 低估 PE%=18.3 PB%=22.1 成交额=39.3亿

⚠️ 可关注清单（数据估算）:
  1. sh513090 证券ETF ⚠️ 低估 PE%=25.3 PB%=28.7 成交额=90.9亿

【数据质量说明】
✅ 真实：真实历史分位（乐咕乐股）
⚠️ 估算：申万行业/穿透估值估算
❓ 未知：数据不可用

【风险提示】本简报仅供参考，不构成投资建议..."""
    
    p = content_frame.paragraphs[0]
    p.text = example_text
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.font.color.rgb = COLOR_DARK
    p.space_after = Pt(0)
    
    # ==================== 第9页：项目成果与下一步计划 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "项目成果与下一步计划"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    text_content = [
        ("已完成", Pt(20), True, COLOR_SUCCESS),
        ("• ✅ 数据真实性保障：仅推荐percentile_real_flag=True的ETF", Pt(16), False, COLOR_DARK),
        ("• ✅ 数据透明化：日报标注数据来源与可信度", Pt(16), False, COLOR_DARK),
        ("• ✅ 穿透估值实现：基于持仓×申万行业PE/PB估算", Pt(16), False, COLOR_DARK),
        ("• ✅ 全流程自动化：01:30→09:00→09:25，09:30前推送", Pt(16), False, COLOR_DARK),
        ("• ✅ Webchat推送修复：4个cron任务已改为webchat窗口推送", Pt(16), False, COLOR_DARK),
        ("", Pt(16), False, COLOR_DARK),
        ("下一步计划", Pt(20), True, COLOR_PRIMARY),
        ("• 🎯 提升真实数据覆盖率（目标80%+）", Pt(16), False, COLOR_DARK),
        ("  - 接入中证指数官网获取官方估值数据", Pt(14), False, COLOR_DARK),
        ("  - 自建历史数据库（ETF自身历史分位）", Pt(14), False, COLOR_DARK),
        ("• 📊 添加回测系统验证策略有效性", Pt(16), False, COLOR_DARK),
        ("• 🔔 优化告警机制（穿透估值失败通知）", Pt(16), False, COLOR_DARK),
    ]
    
    for text, size, bold, color in text_content:
        p = content_frame.add_paragraph() if text else content_frame.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6)
    
    # ==================== 第10页：结束页 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 结束语
    end_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    end_frame = end_box.text_frame
    end_frame.text = "感谢聆听！"
    p = end_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    # 保存
    output_path = "/Users/zhangxianghao/.qclaw/workspace/etf-agent/output/ETF指数智能投资顾问项目汇报.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
