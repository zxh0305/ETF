#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ETF项目进展汇报PPT - 精美完整版（15页）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 配色
C_DARK = RGBColor(31, 78, 121)
C_MED = RGBColor(68, 114, 196)
C_LIGHT = RGBColor(180, 198, 231)
C_ACCENT = RGBColor(237, 125, 49)
C_GREEN = RGBColor(112, 173, 71)
C_BG = RGBColor(242, 242, 242)

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def title_bar(slide, text, color=C_MED):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def box(slide, x, y, w, h, border_color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(245, 248, 255)
    s.line.color.rgb = border_color
    s.line.width = Pt(2)
    return s

def add_text(tf, text, size=14, bold=False, color=None, space=4):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    if bold:
        p.font.bold = True
    if color:
        p.font.color.rgb = color
    p.space_after = Pt(space)
    return p

# ===== Slide 1: 标题页 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C_DARK)
# 顶部条
t = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
t.fill.solid(); t.fill.fore_color.rgb = C_MED; t.line.fill.background()
# 底部条
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6), Inches(10), Inches(1.5))
b.fill.solid(); b.fill.fore_color.rgb = C_MED; b.line.fill.background()
# 标题
tb = s.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.2))
p = tb.text_frame.paragraphs[0]
p.text = "ETF智能投资顾问项目"
p.font.size = Pt(44); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER
# 副标题
tb2 = s.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
p = tb2.text_frame.paragraphs[0]
p.text = "进展汇报"; p.font.size = Pt(36); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER
# 标签
tb3 = s.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.6))
p = tb3.text_frame.paragraphs[0]
p.text = "数据驱动 + 智能体决策  |  本地运行 · 免费开源 · 零代码"
p.font.size = Pt(18); p.font.color.rgb = RGBColor(200, 200, 200); p.alignment = PP_ALIGN.CENTER

# ===== Slide 2: 目录 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C_BG)
title_bar(s, "目录", C_MED)
items = [
    ("01", "ETF与A股市场", C_MED),
    ("02", "项目初衷", C_ACCENT),
    ("03", "项目目标", C_GREEN),
    ("04", "已完成功能", C_DARK),
    ("05", "核心Agent实现", RGBColor(255, 192, 0)),
    ("06", "实际运行效果", C_MED),
    ("07", "当前差距分析", C_ACCENT),
    ("08", "下一步计划", C_GREEN),
    ("09", "总结与展望", C_DARK),
]
y = 1.6
for num, title, color in items:
    nb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.6), Inches(0.5))
    nb.fill.solid(); nb.fill.fore_color.rgb = color; nb.line.fill.background()
    tf = nb.text_frame; tf.vertical_anchor = 1
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(16); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER
    tb = s.shapes.add_textbox(Inches(1.6), Inches(y), Inches(7.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = RGBColor(50, 50, 50)
    y += 0.55

# ===== Slide 3: ETF与A股市场 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, RGBColor(255, 255, 255))
title_bar(s, "01 | ETF与A股市场", C_MED)
# ETF定义框
b1 = box(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(2.7), C_MED)
tf = b1.text_frame; tf.word_wrap = True
add_text(tf, "📈 ETF（交易型开放式指数基金）", 16, True, C_MED, 10)
for t in ["• 在交易所上市交易，可像股票一样买卖", "• 跟踪特定指数，分散投资风险", "• 管理费低（0.15%-0.5%），透明度高", "• 2026年A股ETF数量：超800只"]:
    add_text(tf, t, 12, False, RGBColor(60,60,60), 5)
# 与A股关系框
b2 = box(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.7), C_ACCENT)
tf = b2.text_frame; tf.word_wrap = True
add_text(tf, "📊 ETF与A股市场的关系", 16, True, C_ACCENT, 10)
for t in ["• A股超5000只股票，普通投资者选择困难", "• ETF覆盖宽基、行业、主题、跨境等", "• 投资价值：避免个股风险，获取平均收益", "• 适合量化筛选和自动化交易"]:
    add_text(tf, t, 12, False, RGBColor(60,60,60), 5)
# 底部总结框
b3 = box(s, Inches(0.5), Inches(4.4), Inches(9), Inches(2.7), C_GREEN)
tf = b3.text_frame; tf.word_wrap = True
add_text(tf, "💡 为何选择ETF作为智能化投资标的？", 16, True, C_GREEN, 10)
for t in ["1. 分散风险 - 持有一篮子股票，避免单一个股风险",
          "2. 成本低廉 - 管理费率远低于主动基金（通常0.15%-0.5%）",
          "3. 透明度高 - 持仓、权重、估值实时公开可查",
          "4. 流动性好 - 日成交亿元级别，进出自由",
          "5. 适合自动化 - 标准化的产品特性，便于程序化交易"]:
    add_text(tf, t, 12, False, RGBColor(60,60,60), 5)

# ===== Slide 4: 项目初衷 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, RGBColor(255, 255, 255))
title_bar(s, "02 | 项目初衷", C_MED)
# 三栏
for i, (title, items, color, x) in enumerate([
    ("💡 投资痛点", ["• 人工盯盘耗时耗力，情绪化决策严重", "• ETF数量多（800+只），人工筛选困难", "• 估值数据分散在多个平台", "• 错过最佳买卖时机，收益打折扣", "• 缺乏系统化的投资纪律执行"], C_ACCENT, Inches(0.5)),
    ("🔧 技术契机", ["• OpenClaw（龙虾）智能体开源，自主执行", "• 免费金融数据接口成熟（AkShare）", "• 量化投资从机构走向个人", "• AI技术普及，降低自动化门槛", "• 本地大模型（Ollama）性能提升"], C_GREEN, Inches(3.5)),
    ("🎯 个人需求", ["• 想要24小时在线的'赛博员工'", "• 自动监控市场，推送低估机会", "• 数据本地运行，确保隐私安全", "• 零成本使用，适合个人投资者", "• 系统化执行，避免情绪化决策"], RGBColor(255, 192, 0), Inches(6.5)),
]):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.8), Inches(5.5))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(245, 248, 255)
    b.line.color.rgb = color; b.line.width = Pt(2)
    tf = b.text_frame; tf.word_wrap = True
    add_text(tf, title, 16, True, color, 12)
    for item in items:
        add_text(tf, item, 11, False, RGBColor(60,60,60), 6)

# ===== Slide 5: 项目目标 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C_BG)
title_bar(s, "03 | 项目目标", C_MED)
# 左栏
b1 = box(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), C_MED)
tf = b1.text_frame; tf.word_wrap = True
add_text(tf, "📈 核心功能目标", 18, True, C_MED, 12)
add_text(tf, "1. 数据采集自动化", 14, True, C_DARK, 6)
add_text(tf, "   • 每日09:20自动拉取全市场ETF数据", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 覆盖PE/PB/分位/成交额等核心指标", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 多数据源融合，确保数据完整性", 12, False, RGBColor(60,60,60), 10)
add_text(tf, "2. 智能筛选系统", 14, True, C_DARK, 6)
add_text(tf, "   • 基于PE/PB分位≤30%筛选", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 流动性门槛：日均成交≥1亿元", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • PEG<1确保成长价值匹配", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 评分系统：PE40%+PB30%+流动性30%", 12, False, RGBColor(60,60,60), 10)
add_text(tf, "3. 自动推送与提醒", 14, True, C_DARK, 6)
add_text(tf, "   • 微信每日09:25接收简报", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 变化信号实时提醒", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 操作建议生成", 12, False, RGBColor(60,60,60), 4)
# 右栏
b2 = box(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5), C_ACCENT)
tf = b2.text_frame; tf.word_wrap = True
add_text(tf, "🎯 技术架构目标", 18, True, C_ACCENT, 12)
add_text(tf, "1. 3个独立Agent分工", 14, True, C_DARK, 6)
add_text(tf, "   • Agent1: ETF_估值数据采集", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • Agent2: ETF_低估筛选", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • Agent3: ETF_提醒推送", 12, False, RGBColor(60,60,60), 10)
add_text(tf, "2. 数据驱动决策", 14, True, C_DARK, 6)
add_text(tf, "   • 多数据源融合（乐咕+申万+CNINFO）", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 历史分位计算（20年数据）", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 实时数据更新（新浪接口）", 12, False, RGBColor(60,60,60), 10)
add_text(tf, "3. 安全与隐私", 14, True, C_DARK, 6)
add_text(tf, "   • 本地运行，不上传云端", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 辅助决策，不自动交易", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 数据加密存储", 12, False, RGBColor(60,60,60), 10)
add_text(tf, "📈 预期效果", 14, True, C_DARK, 6)
add_text(tf, "   • 不错过任何低估机会", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 降低情绪化决策干扰", 12, False, RGBColor(60,60,60), 4)
add_text(tf, "   • 提升投资效率和胜率", 12, False, RGBColor(60,60,60), 4)

# ===== Slide 6: 已完成功能 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, RGBColor(255, 255, 255))
title_bar(s, "04 | 已完成功能（完成度88%）", C_GREEN)
# 三个完成度卡片
for i, (title, items, color, x) in enumerate([
    ("✅ 数据采集层（100%）", ["AkShare接口接入", "每日09:20自动采集", "三数据源架构", "数据覆盖率46.4%"], C_GREEN, Inches(0.5)),
    ("✅ 智能筛选层（70%）", ["PE/PB分位≤30%筛选", "流动性筛选（≥5000万）", "评分系统已实现", "变化检测已实现"], RGBColor(255, 192, 0), Inches(3.5)),
    ("✅ 输出推送层（100%）", ["每日09:25微信推送", "格式按PPT模板优化", "变化信号识别", "操作建议生成"], C_MED, Inches(6.5)),
]):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.8), Inches(2.5))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(245, 248, 255)
    b.line.color.rgb = color; b.line.width = Pt(3)
    tf = b.text_frame; tf.word_wrap = True
    add_text(tf, title, 14, True, color, 10)
    for item in items:
        add_text(tf, "• " + item, 12, False, RGBColor(60,60,60), 5)
# 详细说明
tb = s.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.5), Inches(3))
tf = tb.text_frame; tf.word_wrap = True
add_text(tf, "详细说明", 16, True, C_DARK, 10)
for t in [
    "• AkShare接口接入，每日09:20自动采集全市场ETF数据",
    "• 三数据源架构：乐咕乐股 + 申万行业 + CNINFO，覆盖率46.4%（167/360只ETF）",
    "• PE/PB分位≤30%筛选规则已实现，流动性筛选：当日成交额≥5000万",
    "• 每日09:25自动推送微信，推送格式已按PPT模板优化"
]:
    add_text(tf, t, 12, False, RGBColor(60,60,60), 5)

# ===== Slide 7: Agent1数据采集 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C_BG)
title_bar(s, "05 | 核心Agent1：数据采集", C_DARK)
# 左：数据源
b1 = box(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), C_MED)
tf = b1.text_frame; tf.word_wrap = True
add_text(tf, "📊 三数据源架构", 16, True, C_MED, 10)
add_text(tf, "1. 乐咕乐股（主源）", 13, True, C_DARK, 5)
add_text(tf, "   覆盖宽基、行业、商品、QDII等ETF", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   提供PE/PB实时估值+20年历史分位", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   覆盖132只ETF（36.7%）", 11, False, RGBColor(80,80,80), 8)
add_text(tf, "2. 申万行业（补充源）", 13, True, C_DARK, 5)
add_text(tf, "   31个申万一级行业PE/PB数据", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   通过akshare sw_index_first_info获取", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   覆盖134只ETF（37.2%）", 11, False, RGBColor(80,80,80), 8)
add_text(tf, "3. CNINFO巨潮（兜底源）", 13, True, C_DARK, 5)
add_text(tf, "   深交所/上交所公开数据", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   覆盖主动管理LOF等", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   覆盖195只ETF（54.2%）", 11, False, RGBColor(80,80,80), 3)
# 右：运行机制
b2 = box(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5), C_ACCENT)
tf = b2.text_frame; tf.word_wrap = True
add_text(tf, "⚙️ 运行机制", 16, True, C_ACCENT, 10)
add_text(tf, "触发方式", 13, True, C_ACCENT, 5)
add_text(tf, "   • OpenClaw定时任务（cron）", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   • 每个交易日09:20自动触发", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   • 周一至周五（1-5）", 11, False, RGBColor(80,80,80), 8)
add_text(tf, "执行流程", 13, True, C_ACCENT, 5)
add_text(tf, "   1. 拉取全市场ETF列表（800+只）", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   2. 匹配乐咕乐股估值数据", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   3. 申万行业PE/PB补充", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   4. CNINFO兜底补充", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   5. 计算历史分位（20年）", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   6. 输出JSON到latest文件", 11, False, RGBColor(80,80,80), 8)
add_text(tf, "输出文件", 13, True, C_ACCENT, 5)
add_text(tf, "   • etf_valuation_latest.json", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "   • 包含360只有效ETF数据", 11, False, RGBColor(80,80,80), 3)

# ===== Slide 8: Agent2筛选 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C_BG)
title_bar(s, "05 | 核心Agent2：低估筛选", RGBColor(255, 192, 0))
# 左：筛选逻辑
b1 = box(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), RGBColor(255, 192, 0))
tf = b1.text_frame; tf.word_wrap = True
add_text(tf, "📊 筛选逻辑详解", 16, True, RGBColor(255, 192, 0), 10)
add_text(tf, "正式低估池（三个条件同时满足）", 13, True, C_DARK, 6)
add_text(tf, "   ① PE分位 ≤ 30%（近20年历史低位）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   ② PB分位 ≤ 30%（近20年历史低位）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   ③ 当日成交额 ≥ 5000万元", 11, False, RGBColor(80,80,80), 10)
add_text(tf, "观察池（满足任一条件）", 13, True, C_DARK, 6)
add_text(tf, "   ① PE分位 ≤ 30% 或 PB分位 ≤ 30%", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   ② 成交额 ≥ 1000万元（流动性尚可）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   ③ 指数估值偏低（申万行业分位≤20%）", 11, False, RGBColor(80,80,80), 10)
add_text(tf, "评分系统", 13, True, C_DARK, 6)
add_text(tf, "   PE分位权重：40%", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   PB分位权重：30%", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   流动性权重：30%", 11, False, RGBColor(80,80,80), 4)
# 右：运行结果
b2 = box(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5), C_MED)
tf = b2.text_frame; tf.word_wrap = True
add_text(tf, "📈 运行结果（2026-05-11）", 16, True, C_MED, 10)
add_text(tf, "输入数据", 13, True, C_DARK, 6)
add_text(tf, "   • ETF总数：360只", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • 有PE分位数据：167只（46.4%）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • 有PB分位数据：167只（46.4%）", 11, False, RGBColor(80,80,80), 10)
add_text(tf, "筛选结果", 13, True, C_DARK, 6)
add_text(tf, "   • 正式低估池：1只（白酒基金LOF）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "     - PE分位：21.5%", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "     - PB分位：21.7%", 11, False, RGBColor(80,80,80), 3)
add_text(tf, "     - 成交额：0.68亿", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • 观察池：359只（LOF为主）", 11, False, RGBColor(80,80,80), 10)
add_text(tf, "分析说明", 13, True, C_DARK, 6)
add_text(tf, "   • 低估ETF大多为LOF基金", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • LOF流动性普遍较差", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • 需提高成交额门槛（待优化）", 11, False, RGBColor(80,80,80), 4)

# ===== Slide 9: Agent3推送 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C_BG)
title_bar(s, "05 | 核心Agent3：提醒推送", C_GREEN)
# 左：推送内容
b1 = box(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), C_GREEN)
tf = b1.text_frame; tf.word_wrap = True
add_text(tf, "📱 推送内容模板", 16, True, C_GREEN, 10)
add_text(tf, "【ETF低估智能体 | 每日09:20播报】", 12, True, C_DARK, 6)
add_text(tf, "日期：YYYY-MM-DD", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "", 8, False, RGBColor(80,80,80), 0)
add_text(tf, "今日符合条件共：N只", 12, True, C_DARK, 6)
add_text(tf, "", 8, False, RGBColor(80,80,80), 0)
add_text(tf, "重点关注：", 12, True, C_DARK, 6)
add_text(tf, "1. 白酒基金LOF（sz161725）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   PE分位21.5% PB分位21.7% 成交0.68亿", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "", 8, False, RGBColor(80,80,80), 0)
add_text(tf, "变化提醒：", 12, True, C_DARK, 6)
add_text(tf, "🆕 白酒基金LOF新入选低估池", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "📤 嘉实原油LOF移出观察池", 11, False, RGBColor(80,80,80), 4)
# 右：运行配置
b2 = box(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5), C_MED)
tf = b2.text_frame; tf.word_wrap = True
add_text(tf, "⚙️ 运行配置", 16, True, C_MED, 10)
add_text(tf, "定时任务", 13, True, C_DARK, 6)
add_text(tf, "   • 任务名：ETF每日简报推送", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • 触发时间：每日09:25（周一至周五）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • Session Target：isolated", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • Payload：agentTurn", 11, False, RGBColor(80,80,80), 10)
add_text(tf, "推送渠道", 13, True, C_DARK, 6)
add_text(tf, "   • openclaw-weixin（微信插件）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • 直接推送给用户", 11, False, RGBColor(80,80,80), 10)
add_text(tf, "推送格式", 13, True, C_DARK, 6)
add_text(tf, "   • 结构化文本（Markdown格式）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • 包含：符合条件、重点关注、", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "     变化提醒、操作建议、风险提示", 11, False, RGBColor(80,80,80), 10)
add_text(tf, "健康检查", 13, True, C_DARK, 6)
add_text(tf, "   • 每个交易日10:05自动检查", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • 监控Pipeline文件完整性", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   • 失败告警：连续2次失败触发", 11, False, RGBColor(80,80,80), 4)

# ===== Slide 10: 实际运行效果 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, RGBColor(255, 255, 255))
title_bar(s, "06 | 实际运行效果", C_MED)
# 推送示例框
b = box(s, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), C_GREEN)
tf = b.text_frame; tf.word_wrap = True
add_text(tf, "📱 微信推送示例（2026-05-11 实际运行数据）", 18, True, C_GREEN, 15)
add_text(tf, "━" * 60, 10, False, RGBColor(180, 180, 180), 5)
add_text(tf, "【ETF低估智能体 | 每日09:20播报】", 14, True, C_DARK, 8)
add_text(tf, "日期：2026-05-11（周一）", 12, False, RGBColor(80,80,80), 6)
add_text(tf, "", 6, False, RGBColor(80,80,80), 0)
add_text(tf, "今日符合条件共：1只", 14, True, C_DARK, 8)
add_text(tf, "", 6, False, RGBColor(80,80,80), 0)
add_text(tf, "【重点关注】", 13, True, C_ACCENT, 6)
add_text(tf, "1. 白酒基金LOF（sz161725）", 12, True, C_DARK, 5)
add_text(tf, "   PE分位：21.5%  |  PB分位：21.7%  |  成交额：0.68亿", 11, False, RGBColor(80,80,80), 10)
add_text(tf, "", 6, False, RGBColor(80,80,80), 0)
add_text(tf, "【变化提醒】", 13, True, C_ACCENT, 6)
add_text(tf, "🆕 白酒基金LOF新入选低估池", 12, False, RGBColor(80,80,80), 5)
add_text(tf, "📤 嘉实原油LOF移出观察池", 12, False, RGBColor(80,80,80), 10)
add_text(tf, "", 6, False, RGBColor(80,80,80), 0)
add_text(tf, "【操作建议】", 13, True, C_ACCENT, 6)
add_text(tf, "白酒基金LOF估值偏低，建议等待确认后关注。", 12, False, RGBColor(80,80,80), 10)
add_text(tf, "━" * 60, 10, False, RGBColor(180, 180, 180), 5)
add_text(tf, "⚠️ 风险提示：仅供参考，不构成投资建议。", 11, False, RGBColor(150, 150, 150), 5)

# ===== Slide 11: 差距分析 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, RGBColor(255, 255, 255))
title_bar(s, "07 | 当前差距分析（完成度88%）", C_ACCENT)
# 左：主要差距
b1 = box(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), C_ACCENT)
tf = b1.text_frame; tf.word_wrap = True
add_text(tf, "❌ 主要差距", 16, True, C_ACCENT, 10)
add_text(tf, "1. PEG指标未实现", 14, True, C_DARK, 6)
add_text(tf, "   PPT要求：PEG<1（成长价值匹配）", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   当前状态：完全没有PEG数据", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   影响：可能选中价值陷阱", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   解决方案：注册Tushare Pro获取", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   盈利增长率数据", 11, False, RGBColor(80,80,80), 10)
add_text(tf, "2. 20日成交额均值未实现", 14, True, C_DARK, 6)
add_text(tf, "   PPT要求：近20日日均≥1亿元", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   当前状态：使用当日成交额，阈值5000万", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   影响：筛选结果不稳定", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   解决方案：获取近20日历史数据，", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   计算均值，调整阈值为1亿元", 11, False, RGBColor(80,80,80), 4)
# 右：完成度
b2 = box(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5), C_GREEN)
tf = b2.text_frame; tf.word_wrap = True
add_text(tf, "✅ 已完成功能", 16, True, C_GREEN, 10)
for t in ["• 数据采集自动化（09:20定时）", "• PE/PB分位筛选（≤30%）", "• 微信推送（09:25定时）", "• 推送模板（符合PPT格式）", "• 定时任务+健康检查+告警"]:
    add_text(tf, t, 12, False, RGBColor(80,80,80), 6)
add_text(tf, "", 8, False, RGBColor(80,80,80), 0)
add_text(tf, "📊 完成度评估", 14, True, C_DARK, 8)
add_text(tf, "   Agent1: 数据采集  100%", 12, False, RGBColor(80,80,80), 4)
add_text(tf, "   Agent2: 低估筛选   70%（缺PEG）", 12, False, RGBColor(80,80,80), 4)
add_text(tf, "   Agent3: 提醒推送  100%", 12, False, RGBColor(80,80,80), 4)
add_text(tf, "   推送格式           100%", 12, False, RGBColor(80,80,80), 8)
add_text(tf, "━━━━━━━━━━━━━━━━━━━━", 10, False, RGBColor(180,180,180), 5)
add_text(tf, "🎯 总体完成度：88%", 16, True, C_GREEN, 5)

# ===== Slide 12: 下一步计划 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C_BG)
title_bar(s, "08 | 下一步计划", C_GREEN)
# 优先级卡片
for i, (title, items, color, x) in enumerate([
    ("🔴 高优先级（Week 1）", ["实现PEG指标（2-3天）", "• 注册Tushare Pro获取盈利增长率", "• 计算PEG = PE / 盈利增长率", "• 增加筛选条件：PEG<1", "", "实现20日成交额均值（1天）", "• 获取近20日历史成交额数据", "• 计算均值，调整阈值为1亿"], C_ACCENT, Inches(0.5)),
    ("🟡 中优先级（Week 2）", ["风险评估模块（1-2天）", "• 基于波动率、最大回撤评估风险", "• 生成风险等级标签（高/中/低）", "• 纳入筛选考量因素", "", "配置建议模块（2-3天）", "• 基于低估程度、流动性推荐仓位", "• 给出建仓/加仓/观望建议"], RGBColor(255, 192, 0), Inches(3.5)),
    ("🟢 低优先级（Week 3+）", ["可视化报告生成器", "• 生成PDF/HTML格式报告", "• 包含图表、数据可视化", "• 支持一键导出", "", "Web界面部署（可选）", "• 搭建简易Web dashboard", "• 实时查看低估池变化"], C_GREEN, Inches(6.5)),
]):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.8), Inches(5.5))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(245, 248, 255)
    b.line.color.rgb = color; b.line.width = Pt(2)
    tf = b.text_frame; tf.word_wrap = True
    add_text(tf, title, 14, True, color, 10)
    for item in items:
        if item.startswith("实现") or item.startswith("风险") or item.startswith("配置") or item.startswith("可视") or item.startswith("Web"):
            add_text(tf, item, 12, True, C_DARK, 6)
        else:
            add_text(tf, item, 11, False, RGBColor(80,80,80), 4)

# ===== Slide 13: 总结 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, RGBColor(255, 255, 255))
title_bar(s, "09 | 总结与展望", C_DARK)
# 左：总结
b1 = box(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), C_MED)
tf = b1.text_frame; tf.word_wrap = True
add_text(tf, "📊 项目总结", 16, True, C_MED, 10)
for t in ["• 核心功能已基本完成（88%）", "• 数据采集→筛选→推送全流程已打通", "• 定时任务稳定运行，微信推送正常"]:
    add_text(tf, t, 12, False, RGBColor(80,80,80), 6)
add_text(tf, "", 8, False, RGBColor(80,80,80), 0)
add_text(tf, "💡 核心亮点", 14, True, C_DARK, 8)
for t in ["• 数据驱动：多数据源融合，覆盖46.4%", "• 智能体决策：OpenClaw自主执行", "• 本地安全运行：隐私绝对保障", "• 免费开源：零成本使用"]:
    add_text(tf, t, 12, False, RGBColor(80,80,80), 6)
add_text(tf, "", 8, False, RGBColor(80,80,80), 0)
add_text(tf, "⚠️ 当前限制", 14, True, C_DARK, 8)
for t in ["• PEG指标缺失（需Tushare Pro）", "• 20日成交额均值未实现", "• 部分ETF无PE/PB数据（主动LOF）"]:
    add_text(tf, t, 12, False, RGBColor(80,80,80), 6)
# 右：展望
b2 = box(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5), C_ACCENT)
tf = b2.text_frame; tf.word_wrap = True
add_text(tf, "🚀 未来展望", 16, True, C_ACCENT, 10)
add_text(tf, "1. AI智能预测", 13, True, C_DARK, 6)
add_text(tf, "   引入机器学习模型，预测ETF未来走势", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   提升筛选胜率", 11, False, RGBColor(80,80,80), 8)
add_text(tf, "2. 多因子筛选模型", 13, True, C_DARK, 6)
add_text(tf, "   构建更复杂的多因子模型", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   多维度评估资产，提升准确性", 11, False, RGBColor(80,80,80), 8)
add_text(tf, "3. 策略分享社区", 13, True, C_DARK, 6)
add_text(tf, "   建立开放社区，允许用户分享策略", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   实现智慧共享", 11, False, RGBColor(80,80,80), 8)
add_text(tf, "4. 实盘对接（可选）", 13, True, C_DARK, 6)
add_text(tf, "   模拟盘测试验证策略有效性", 11, False, RGBColor(80,80,80), 4)
add_text(tf, "   用户自主选择是否实盘", 11, False, RGBColor(80,80,80), 4)

# ===== Slide 14: 结束页 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C_DARK)
# 顶部装饰
t = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
t.fill.solid(); t.fill.fore_color.rgb = C_MED; t.line.fill.background()
# 底部装饰
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6), Inches(10), Inches(1.5))
b.fill.solid(); b.fill.fore_color.rgb = C_MED; b.line.fill.background()
# 标题
tb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
p = tb.text_frame.paragraphs[0]
p.text = "感谢聆听"; p.font.size = Pt(44); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER
# 副标题
tb2 = s.shapes.add_textbox(Inches(1), Inches(3.6), Inches(8), Inches(1.5))
p = tb2.text_frame.paragraphs[0]
p.text = "ETF智能投资顾问项目进展汇报\n\n本地安全 · 免费 · 零代码\n\n祝你投资顺利，风控为先"
p.font.size = Pt(18); p.font.color.rgb = RGBColor(200, 200, 200); p.alignment = PP_ALIGN.CENTER

# 保存
out = "/Users/zhangxianghao/.qclaw/workspace/etf-agent/output/ETF智能投资顾问项目进展汇报_精美完整版.pptx"
prs.save(out)
print(f"✅ 精美完整版PPT已生成！")
print(f"   文件: {out}")
print(f"   总页数: {len(prs.slides)} 页")
