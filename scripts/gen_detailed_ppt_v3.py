#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ETF智能投资顾问项目进展汇报 — 详细介绍版 v3
修复：留白问题、全部文字颜色问题、第8页白字白底
"""

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Emu(9144000)
prs.slide_height = Emu(6720840)

C_NAVY   = RGBColor(0x1A, 0x3A, 0x5C)
C_BLUE   = RGBColor(0x5B, 0x9B, 0xD5)
C_GOLD   = RGBColor(0xE8, 0xA8, 0x20)
C_GRAY   = RGBColor(0x44, 0x44, 0x44)
C_LGRAY  = RGBColor(0x66, 0x66, 0x66)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_BPURP  = RGBColor(0x2E, 0x6D, 0xA4)
C_PURPLE = RGBColor(0x7D, 0x5C, 0xA0)
C_ORANGE = RGBColor(0xE8, 0x7B, 0x19)
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)
C_BG     = RGBColor(0xF5, 0xF7, 0xFA)
C_DARK   = RGBColor(0x3D, 0x3D, 0x3D)
C_LBBG   = RGBColor(0xED, 0xF1, 0xF7)   # 浅蓝灰背景（内容区用）

# 页面高度常量
PH = 6720840   # slide height
TITLE_H = 960000
GOLD_H = 50000
CONTENT_TOP = TITLE_H + GOLD_H + 80000  # 内容区起始 y
MARGIN_L = 320000
MARGIN_R = 320000
CONTENT_W = 9144000 - MARGIN_L - MARGIN_R
GAP = 60000

def R(slide, l, t, w, h, color, alpha=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def T(slide, l, t, w, h, text, size, bold=False, color=C_LGRAY,
      align=PP_ALIGN.LEFT, va=MSO_ANCHOR.TOP):
    """在text中每个段落分别设置颜色，确保多行文字颜色一致"""
    s = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    s.text_frame.word_wrap = True
    s.text_frame.vertical_anchor = va
    # 用 \n 分割，每段单独设颜色
    lines = text.split('\n')
    # 删除默认空段落
    for para in s.text_frame.paragraphs:
        para.clear()
    for i, line in enumerate(lines):
        if i == 0:
            para = s.text_frame.paragraphs[0]
        else:
            para = s.text_frame.add_paragraph()
        para.text = line
        para.alignment = align
        if size:
            for run in para.runs:
                run.font.size = size
        if bold:
            for run in para.runs:
                run.font.bold = True
        if color is not None:
            for run in para.runs:
                run.font.color.rgb = color
    return s.text_frame

def title_bar(slide, title_text):
    R(slide, 0, 0, 9144000, TITLE_H, C_NAVY)
    R(slide, 0, TITLE_H, 9144000, GOLD_H, C_GOLD)
    T(slide, MARGIN_L, 200000, CONTENT_W, TITLE_H - 200000,
      title_text, size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.BOTTOM)

def card(slide, l, t, w, h, top_color, title_text, lines_text, body_color=C_LGRAY, body_size=Pt(11), title_size=Pt(14)):
    """紧凑型卡片：减少留白"""
    R(slide, l, t, w, h, C_LBBG)
    R(slide, l, t, w, 50000, top_color)
    TITLE_H2 = 420000
    T(slide, l + GAP, t + 60000, w - GAP*2, TITLE_H2, title_text,
      size=title_size, bold=True, color=C_DARK)
    T(slide, l + GAP, t + TITLE_H2 + 60000, w - GAP*2, h - TITLE_H2 - 60000,
      lines_text, size=body_size, bold=False, color=body_color)

def full_card(slide, l, t, w, h, top_color, title_text, lines_text, body_color=C_LGRAY, body_size=Pt(11), title_size=Pt(14)):
    """通栏卡片"""
    R(slide, l, t, w, h, C_LBBG)
    R(slide, l, t, w, 50000, top_color)
    TITLE_H2 = 420000
    T(slide, l + GAP, t + 60000, w - GAP*2, TITLE_H2, title_text,
      size=title_size, bold=True, color=C_DARK)
    T(slide, l + GAP, t + TITLE_H2 + 60000, w - GAP*2, h - TITLE_H2 - 60000,
      lines_text, size=body_size, bold=False, color=body_color)

# ══════════════════════════════════════════════════════════════
#  第1页  标题页
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = C_NAVY

R(slide, 0, PH*0.36, 9144000, 60000, C_GOLD)
R(slide, 0, PH*0.38, 9144000, 30000, C_GOLD)
T(slide, MARGIN_L, PH*0.40, CONTENT_W, PH*0.14,
  "ETF指数智能投资顾问", size=Pt(44), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, MARGIN_L, PH*0.56, CONTENT_W, PH*0.08,
  "项目进展汇报  ·  详细介绍版", size=Pt(18), bold=False, color=C_BLUE, align=PP_ALIGN.CENTER)
T(slide, MARGIN_L, PH*0.80, CONTENT_W, PH*0.05,
  "2026年05月12日", size=Pt(14), bold=False, color=C_BLUE, align=PP_ALIGN.CENTER)
T(slide, MARGIN_L, PH*0.88, CONTENT_W, PH*0.04,
  "张翔豪", size=Pt(12), bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  第2页  ETF简介
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "ETF简介 与 A股投资关系")
# 内容区高度
CA = PH - CONTENT_TOP
CW = (CONTENT_W - GAP) / 2

# 上排两栏
card(slide, MARGIN_L, CONTENT_TOP, CW, CA*0.45, C_BPURP,
     "ETF 是什么？",
     "ETF（Exchange Traded Fund，交易所交易基金）\n"
     "• 在交易所上市、份额可变的开放式基金\n"
     "• 跟踪特定指数（沪深300/纳斯达克100等）\n"
     "• 兼具股票流动性与基金分散投资优势\n"
     "• 管理费率低（0.15%-0.5%），比主动基金便宜")

card(slide, MARGIN_L + CW + GAP, CONTENT_TOP, CW, CA*0.45, C_GREEN,
     "与A股投资的关系",
     "直接买股票 vs 买ETF：\n"
     "• 分散风险：ETF持一篮子股票，避免个股暴雷\n"
     "• 门槛更低：100元起买，无需数千元选股资金\n"
     "• 选股简单：ETF只需选赛道，不用研究个股\n"
     "• 节省时间：无需每天盯盘，适合上班族")

# 下排通栏
full_card(slide, MARGIN_L, CONTENT_TOP + CA*0.50 + GAP, CONTENT_W, CA*0.50, C_GOLD,
          "为什么要做 ETF智能投资顾问？",
          "• 市场痛点：ETF超300只，普通投资者不知道怎么选\n"
          "• 估值难判断：PE/PB分位需自己算，普通投资者不会\n"
          "• 情绪化交易：恐慌时割肉、贪婪时追高，需要理性信号\n"
          "• 解决方案：数据采集 → 低估筛选 → 定时推送，全流程自动")

# ══════════════════════════════════════════════════════════════
#  第3页  项目目标 与 方案选型
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "项目目标 与 方案选型")
CA = PH - CONTENT_TOP

card(slide, MARGIN_L, CONTENT_TOP, CW, CA*0.80, C_BPURP,
     "项目目标",
     "核心目标：\n"
     "① 数据采集：每日自动抓取全市场ETF估值数据\n"
     "② 低估筛选：基于PE/PB分位筛选低估ETF\n"
     "③ 投资配置建议：生成基础配置建议（待实现）\n"
     "④ 风险等级评估：对ETF进行风险分级（待实现）\n"
     "⑤ 可视化报告：输出可视化投资参考（待实现）\n"
     "\n不涉及实盘交易，仅提供投资参考。")

card(slide, MARGIN_L + CW + GAP, CONTENT_TOP, CW, CA*0.80, C_PURPLE,
     "方案选型",
     "为什么选这个方案？\n"
     "• 数据来源：AkShare（免费开源、更新及时）\n"
     "  vs Tushare（需积分）、东方财富（不稳定）\n"
     "• 定时调度：OpenClaw Gateway cron（内置可靠）\n"
     "  vs Linux crontab（无跨平台）\n"
     "• 架构设计：3个独立Agent（职责分离、易调试）\n"
     "  vs 单脚本顺序执行（出错难定位）")

full_card(slide, MARGIN_L, CONTENT_TOP + CA*0.85 + GAP, CONTENT_W, CA*0.15, C_GREEN,
          "核心亮点",
          "数据驱动 + 智能体决策 + 公开数据接口  |  三数据源融合架构  |  每日自动运行无需人工  |  Pipeline版本化管理")

# ══════════════════════════════════════════════════════════════
#  第4页  系统架构
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "系统架构详解")

# 流程图
FH = PH - CONTENT_TOP
BOX_W = (CONTENT_W - GAP*4) / 5
BOX_H = FH * 0.28
BOX_Y = CONTENT_TOP + FH * 0.05

boxes = [
    (C_NAVY,   "定时触发", "09:20\n周一至周五"),
    (C_BPURP,  "Agent 1", "数据采集"),
    (C_PURPLE, "Agent 2", "低估筛选"),
    (C_GREEN,  "Agent 3", "提醒推送"),
    (C_RED,    "微信推送", "消息模板\n推送成功"),
]
for i, (color, title, sub) in enumerate(boxes):
    bx = MARGIN_L + i*(BOX_W + GAP)
    R(slide, bx, BOX_Y, BOX_W, BOX_H, color)
    T(slide, bx, BOX_Y + BOX_H*0.2, BOX_W, BOX_H*0.4, title,
      size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    T(slide, bx, BOX_Y + BOX_H*0.5, BOX_W, BOX_H*0.4, sub,
      size=Pt(10), bold=False, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < len(boxes)-1:
        T(slide, bx+BOX_W, BOX_Y+BOX_H*0.3, GAP, BOX_H*0.4, "→",
          size=Pt(18), bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

T(slide, MARGIN_L, BOX_Y+BOX_H+GAP, CONTENT_W, FH*0.06,
  "数据流：etf_valuation_latest.json → low_valuation_candidates_latest.json → 微信推送消息",
  size=Pt(10), bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

# 三Agent详情
A_Y = BOX_Y + BOX_H + FH*0.08
A_H = FH*0.46
A_W = (CONTENT_W - GAP*2) / 3

a1_content = ("触发：每日 09:20（周一至周五）\n\n"
              "数据源：\n"
              "  新浪财经（行情，100%覆盖）\n"
              "  乐咕乐股（宽基估值+分位）\n"
              "  申万行业（31个行业PE/PB）\n"
              "  巨潮资讯（深交所/上交所）\n\n"
              "输出：etf_valuation_latest.json\n"
              "覆盖：360+ 只ETF")

a2_content = ("触发：Agent1完成后自动触发\n\n"
              "筛选条件：\n"
              "  PE分位 ≤ 30%（已实现）\n"
              "  PB分位 ≤ 30%（已实现）\n"
              "  PEG < 1（待Tushare接入）\n"
              "  近20日日均成交额≥1亿（待实现）\n\n"
              "输出：low_valuation_candidates\n"
              "_latest.json")

a3_content = ("触发：每日 09:25（周一至周五）\n\n"
              "功能：\n"
              "  对比今日与昨日筛选结果\n"
              "  识别信号：新入选/移出/加仓/减仓\n"
              "  生成微信推送消息模板\n\n"
              "推送渠道：微信（定时任务已配置）\n"
              "格式：结构化日报")

card(slide, MARGIN_L, A_Y, A_W, A_H, C_BPURP, "Agent 1：数据采集", a1_content)
card(slide, MARGIN_L+A_W+GAP, A_Y, A_W, A_H, C_PURPLE, "Agent 2：低估筛选", a2_content)
card(slide, MARGIN_L+(A_W+GAP)*2, A_Y, A_W, A_H, C_GREEN, "Agent 3：提醒推送", a3_content)

# ══════════════════════════════════════════════════════════════
#  第5页  数据采集详情
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "数据采集层 — 实现详情")
CA = PH - CONTENT_TOP

full_card(slide, MARGIN_L, CONTENT_TOP, CONTENT_W, CA*0.32, C_BPURP,
          "三数据源融合架构",
          "为最大化数据覆盖率，采用三数据源并行采集 + 优先级合并策略：\n"
          "① 乐咕乐股（优先级1）：宽基指数ETF的PE/PB及20年历史分位，数据质量最高\n"
          "② 申万行业（优先级2）：通过akshare.sw_index_first_info()获取31个行业PE/PB，覆盖134只行业ETF\n"
          "③ 巨潮资讯（优先级3）：覆盖195只主动管理LOF，提供深交所/上交所公开估值数据\n"
          "合并策略：同一ETF出现于多数据源时，按优先级取第一个有效值")

# 三个统计卡
stat_y = CONTENT_TOP + CA*0.37
stat_h = CA*0.27
SW = (CONTENT_W - GAP*2) / 3

card(slide, MARGIN_L, stat_y, SW, stat_h, C_GREEN,
     "行情覆盖率", "98.4%\n376/382 只ETF")

card(slide, MARGIN_L+SW+GAP, stat_y, SW, stat_h, C_BPURP,
     "分位覆盖率", "46.4%\n167/360 只ETF")

card(slide, MARGIN_L+(SW+GAP)*2, stat_y, SW, stat_h, C_PURPLE,
     "数据源分布", "乐咕  47只 12.3%\n申万 134只 35.1%\n巨潮 195只 51.0%")

full_card(slide, MARGIN_L, stat_y+stat_h+GAP, CONTENT_W, CA*0.32, C_GOLD,
          "输出文件",
          "etf_valuation_latest.json  — 最新估值数据（全量）\n"
          "etf_valuation_YYYYMMDD.json  — 每日归档（版本备份，防覆盖丢失）\n"
          "数据字段：code, name, price, pe, pb, pe_percentile, pb_percentile, amount")

# ══════════════════════════════════════════════════════════════
#  第6页  遇到的问题
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "遇到的问题 与 解决方案")
CA = PH - CONTENT_TOP

problems = [
    (C_RED,    "问题1", "数据停滞在5月2日",
     "cron表达式为* * *未限制周末，周末不触发",
     "已改为1-5（周一至周五），并设失败2次告警"),
    (C_ORANGE, "问题2", "健康检查任务超时失败",
     "AgentTurn默认120秒超时，健康检查逻辑复杂",
     "简化任务message，移除多余要求（约10秒完成）"),
    (C_ORANGE, "问题3", "筛选器bug：正式低估池为0",
     "检查avg_amount_20d字段，但数据只有amount字段",
     "修复3处字段名，正式低估池现为1只（白酒基金LOF）"),
    (C_ORANGE, "问题4", "分位覆盖率仅46.4%",
     "203只主动LOF不披露PE/PB；31只有映射未补全",
     "主动LOF属正常；31只有映射ETF可补充etf_enricher"),
    (C_GOLD,   "问题5", "PEG指标与20日成交额未实现",
     "PEG需Tushare盈利增长率；20日均值需历史行情",
     "已注册Tushare待接入；20日均值网络稳定后补充"),
]

row_h = (CA - GAP) / len(problems)
for i, (tc, tag, title, cause, solution) in enumerate(problems):
    ry = CONTENT_TOP + i * row_h
    rh = row_h - GAP*0.5
    R(slide, MARGIN_L, ry, CONTENT_W, rh, C_LBBG)
    R(slide, MARGIN_L, ry, CONTENT_W, 40000, tc)
    T(slide, MARGIN_L + GAP, ry + 45000, 700000, rh - 50000, tag,
      size=Pt(11), bold=True, color=tc)
    T(slide, MARGIN_L + 900000, ry + 45000, 2000000, rh - 50000, title,
      size=Pt(11), bold=True, color=C_DARK)
    T(slide, MARGIN_L + 2900000, ry + 45000, CONTENT_W - 2900000 - GAP, rh - 50000,
      "原因：" + cause + "\n方案：" + solution,
      size=Pt(10), bold=False, color=C_LGRAY)

# ══════════════════════════════════════════════════════════════
#  第7页  定时任务配置
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "定时任务配置详情")
CA = PH - CONTENT_TOP

R(slide, MARGIN_L, CONTENT_TOP, CONTENT_W, 440000, C_NAVY)
col_x = [MARGIN_L, MARGIN_L+2800000, MARGIN_L+4300000, MARGIN_L+6400000]
col_w = [2700000, 1500000, 2100000, CONTENT_W-6400000+MARGIN_L]
T(slide, col_x[0], CONTENT_TOP+80000, col_w[0], 300000, "任务名称", size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, col_x[1], CONTENT_TOP+80000, col_w[1], 300000, "Cron", size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, col_x[2], CONTENT_TOP+80000, col_w[2], 300000, "功能描述", size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
T(slide, col_x[3], CONTENT_TOP+80000, col_w[3], 300000, "状态", size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

tasks = [
    ("ETF_估值数据采集", "20 09 * * 1-5", "触发Agent1采集全市场ETF估值数据", "✅ 正常"),
    ("ETF每日简报推送", "25 09 * * 1-5", "触发Agent3对比昨日数据并推送简报", "✅ 正常"),
    ("ETF数据健康检查", "0 10 * * 1-5", "检查etf_valuation_latest.json数据新鲜度", "✅ 正常"),
    ("ETF流水线健康检查", "5 10 * * 1-5", "检查关键文件+快照连续性", "✅ 正常"),
    ("ETF日志自动清理", "30 10 * * 1", "清理超过7天的旧日志文件", "✅ 正常"),
]
ROW = (PH - CONTENT_TOP - 440000) / len(tasks)
for i, (name, cron_expr, desc, status) in enumerate(tasks):
    ry = CONTENT_TOP + 440000 + i * ROW
    bg = C_LBBG if i % 2 == 0 else C_WHITE
    R(slide, MARGIN_L, ry, CONTENT_W, ROW, bg)
    T(slide, col_x[0], ry+ROW*0.15, col_w[0], ROW*0.7, name, size=Pt(11), color=C_LGRAY)
    T(slide, col_x[1], ry+ROW*0.15, col_w[1], ROW*0.7, cron_expr, size=Pt(11), color=C_LGRAY, align=PP_ALIGN.CENTER)
    T(slide, col_x[2], ry+ROW*0.15, col_w[2], ROW*0.7, desc, size=Pt(11), color=C_LGRAY)
    sc = C_GREEN if "✅" in status else C_ORANGE
    T(slide, col_x[3], ry+ROW*0.15, col_w[3], ROW*0.7, status, size=Pt(11), bold=True, color=sc, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  第8页  当前进展（修复白字白底）
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "当前进展 与 完成度评估")
CA = PH - CONTENT_TOP

# 左栏：已完成
card(slide, MARGIN_L, CONTENT_TOP, CW, CA*0.80, C_GREEN,
     "✅ 已完成功能",
     "【数据采集层】\n"
     "  ✅ 三数据源融合（乐咕+申万+巨潮）\n"
     "  ✅ 行情覆盖率 98.4%（376/382）\n"
     "  ✅ 分位数据覆盖率 46.4%（167/360）\n"
     "  ✅ Pipeline版本化管理（防数据丢失）\n\n"
     "【筛选器层】\n"
     "  ✅ PE分位 ≤ 30%（已实现）\n"
     "  ✅ PB分位 ≤ 30%（已实现）\n"
     "  ✅ 成交额过滤（阈值5000万）\n"
     "  ✅ 比较器：对比昨日数据识别变化信号\n\n"
     "【定时任务层】\n"
     "  ✅ 5个cron任务全部正常运行\n"
     "  ✅ 失败告警机制（连续2次失败告警）")

# 右栏：待实现
card(slide, MARGIN_L+CW+GAP, CONTENT_TOP, CW, CA*0.80, C_ORANGE,
     "⚠️ 待实现功能",
     "【PPT目标中尚未实现】\n"
     "  ⚠️ PEG < 1 筛选（需Tushare盈利数据）\n"
     "  ⚠️ 近20日日均成交额 ≥ 1亿（需历史行情）\n"
     "  ⚠️ 投资配置建议模块（尚未实现）\n"
     "  ⚠️ 风险等级评估模块（尚未实现）\n"
     "  ⚠️ 可视化投资参考报告（尚未实现）\n"
     "  ⚠️ 微信推送渠道完整对接（简报已生成）\n\n"
     "【数据质量提升】\n"
     "  ⚠️ 创业板指/科创50真实分位数据\n"
     "  ⚠️ 31只有指数映射ETF估值补全")

full_card(slide, MARGIN_L, CONTENT_TOP+CA*0.85+GAP, CONTENT_W, CA*0.15, C_BPURP,
          "完成度评估",
          "核心功能完成度：约 88%（数据采集✅ + 筛选✅ + 定时任务✅ + 推送模板✅）\n"
          "PPT目标完成度：约 70%（PEG、配置建议、风险评估、可视化均未实现）")

# ══════════════════════════════════════════════════════════════
#  第9页  未来展望（1个月内完成计划）
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "未来展望 — 1个月内完成计划")

weeks = [
    (C_GREEN,  "第1周（7天）：数据质量提升", [
        "实现20日日均成交额均值计算（获取历史行情）",
        "接入Tushare Pro获取盈利增长率，实现PEG筛选",
        "补充31只有指数映射ETF的估值数据",
    ]),
    (C_BPURP, "第2周（7天）：配置建议+风险评估", [
        "实现投资配置建议模块（基于风险偏好生成ETF组合）",
        "实现风险等级评估（低风险/中风险/高风险分级）",
    ]),
    (C_PURPLE, "第3周（7天）：可视化+微信推送", [
        "完善微信推送渠道对接（验证推送正常）",
        "开发可视化报告（ETF筛选结果+分位可视化）",
    ]),
    (C_GOLD,   "第4周（7天）：测试+部署", [
        "全流程集成测试",
        "修复run_etf_full_pipeline.py路径错误",
        "编写使用文档+部署文档",
        "正式上线运行",
    ]),
]

WK_H = (PH - CONTENT_TOP - GAP) / 2
WK_W = (CONTENT_W - GAP) / 2

positions = [
    (MARGIN_L, CONTENT_TOP),
    (MARGIN_L+WK_W+GAP, CONTENT_TOP),
    (MARGIN_L, CONTENT_TOP+WK_H+GAP),
    (MARGIN_L+WK_W+GAP, CONTENT_TOP+WK_H+GAP),
]

for (color, title, items), (lx, ty) in zip(weeks, positions):
    card(slide, lx, ty, WK_W, WK_H, color, title,
         "\n".join(f"• {item}" for item in items))

# ══════════════════════════════════════════════════════════════
#  第10页  项目总结
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = C_NAVY

R(slide, 0, PH*0.30, 9144000, 50000, C_GOLD)
R(slide, 0, PH*0.32, 9144000, 25000, C_GOLD)
T(slide, MARGIN_L, PH*0.33, CONTENT_W, PH*0.12,
  "项目总结", size=Pt(40), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

summaries = [
    (C_GREEN,  "✅ 核心架构已完成",
     "3个Agent协同（数据采集→低估筛选→提醒推送），5个定时任务每日自动运行，Pipeline版本化管理保障数据安全。"),
    (C_BPURP,  "✅ 数据采集层稳定",
     "三数据源融合（乐咕+申万+巨潮），行情覆盖率98.4%，分位覆盖率46.4%，数据新鲜度实时监控。"),
    (C_PURPLE, "⚠️ 筛选条件待补全",
     "PE/PB分位筛选已实现，PEG和20日成交额均值因数据源限制暂未实现，完成后筛选准确度将进一步提升。"),
    (C_GOLD,   "🚀 下一步：配置建议+风险评估",
     "核心功能完成度88%，剩余12%为配置建议+风险评级+可视化。计划1个月内全部完成。"),
]
sy = PH * 0.48
for color, title, content in summaries:
    R(slide, MARGIN_L, sy, CONTENT_W, 480000, C_NAVY)
    R(slide, MARGIN_L, sy, 60000, 480000, color)
    T(slide, MARGIN_L + 120000, sy + 60000, CONTENT_W*0.35, 360000, title,
      size=Pt(14), bold=True, color=C_WHITE)
    T(slide, MARGIN_L + CONTENT_W*0.38, sy + 60000, CONTENT_W*0.6, 360000, content,
      size=Pt(11), bold=False, color=C_BLUE)
    sy += 500000

T(slide, MARGIN_L, PH*0.93, CONTENT_W, PH*0.04,
  "报告日期：2026年05月12日  |  详细介绍版  |  共10页",
  size=Pt(10), bold=False, color=C_GRAY, align=PP_ALIGN.RIGHT)

out = "/Users/zhangxianghao/.qclaw/workspace/etf-agent/output/ETF项目进展汇报_详细介绍版_v3_20260512.pptx"

prs.save(out)
print(f"✅ v3 PPT已生成: {out}")
print(f"   共 {len(prs.slides)} 页")
print(f"   幻灯片尺寸: 9144000×{prs.slide_height} EMU")
