#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成ETF智能体优化汇报PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
from pathlib import Path

# ====== 配色（明亮风格） ======
BG_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_GREEN= RGBColor(0x2E, 0xCC, 0x71)
ACCENT_ORANGE=RGBColor(0xF3, 0x9C, 0x12)
ACCENT_RED  = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_BLUE  = RGBColor(0xD6, 0xEA, 0xF8)
LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
LIGHT_ORANGE= RGBColor(0xFD, 0xEB, 0xC7)
LIGHT_GRAY  = RGBColor(0xF2, 0xF3, 0xF4)
MED_GRAY    = RGBColor(0x95, 0xA5, 0xA6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SUBTITLE_GRAY=RGBColor(0x7F, 0x8C, 0x8D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='PingFang SC'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(tf, text, font_size=16, color=DARK_TEXT, bold=False, space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT, font_name='PingFang SC'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p

# ====== 第1页：封面 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白
add_bg(slide, WHITE)

# 顶部装饰条
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2), Inches(10), Inches(1.2),
             'ETF低估智能体 · 两日优化汇报', font_size=40, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.6),
             '数据覆盖率 67.9% → 87.3% | 推送链路全面加固', font_size=22, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
             '2026.05.14 — 2026.05.15', font_size=18, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

# ====== 第2页：优化总览 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             '优化总览', font_size=30, color=DARK_TEXT, bold=True)

# 4个卡片
cards = [
    ('🔧', '推送链路修复', '微信Session过期→静默失败\n扫码登录+重启Gateway修复\n加Session健康检查4次/天', LIGHT_BLUE, ACCENT_BLUE),
    ('⏰', '推送时序修复', '日报09:25→09:40\n流水线~9分钟完成\n确保读取当日最新数据', LIGHT_GREEN, ACCENT_GREEN),
    ('📊', '行业穿透估值', '5854只股票→30个申万行业\n重仓股行业权重×行业PE/PB\n282只ETF获分位数据', LIGHT_ORANGE, ACCENT_ORANGE),
    ('🔄', '流水线集成', '新增step2a穿透合并(0.1秒)\n11步骤9.5分钟完整运行\n正式低估池35→38只', LIGHT_GRAY, ACCENT_BLUE),
]

for i, (icon, title, desc, bg, accent) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.1)
    y = Inches(1.4 + row * 2.8)
    
    card = add_shape(slide, x, y, Inches(5.6), Inches(2.4), bg, accent, Pt(2))
    
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(5), Inches(0.5),
                 f'{icon}  {title}', font_size=22, color=accent, bold=True)
    
    txBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.8), Inches(5), Inches(1.4))
    tf = txBox.text_frame; tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = DARK_TEXT
            tf.paragraphs[0].font.name = 'PingFang SC'
        else:
            add_para(tf, line, font_size=16, color=DARK_TEXT)

# ====== 第3页：推送链路修复详情 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             '🔧 Day1 推送链路修复', font_size=30, color=DARK_TEXT, bold=True)

# 时间线
timeline = [
    ('14:14', '用户反馈未收到推送', 'cron显示delivered但微信无消息'),
    ('14:20', '发现delivery降级到fallback', '微信通道瞬时不可用，系统静默降级'),
    ('14:26', '根因定位：Session过期', 'errcode:-14 session timeout，token失效'),
    ('14:30', '扫码登录+重启Gateway', '新token需重启才加载，用户需先发消息激活'),
    ('14:38', '加防护1：Session健康检查', '每日4次(07/11/15/19)检查，过期告警'),
    ('14:51', '加防护2：主动推送模式', 'cron delivery改none，agent内部调message工具'),
]

for i, (time, title, desc) in enumerate(timeline):
    y = Inches(1.3 + i * 0.95)
    # 时间标签
    add_shape(slide, Inches(0.8), y, Inches(1.3), Inches(0.7), ACCENT_BLUE)
    add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(1.3), Inches(0.5),
                 time, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 内容
    add_text_box(slide, Inches(2.3), y, Inches(4), Inches(0.35),
                 title, font_size=17, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(2.3), y + Inches(0.35), Inches(4), Inches(0.35),
                 desc, font_size=13, color=SUBTITLE_GRAY)

# 右侧教训卡片
card = add_shape(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(5.2), LIGHT_ORANGE, ACCENT_ORANGE, Pt(2))
add_text_box(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.5),
             '💡 关键教训', font_size=20, color=ACCENT_ORANGE, bold=True)

lessons = [
    '微信Bot Session过期后，\nmessage工具仍返回成功（静默失败）',
    '扫码登录后必须重启Gateway，\n运行中的进程不会加载新token',
    'cron的delivery模式降级时无人感知，\n需改用主动推送才能发现错误',
    'session-guard暂停账号1小时，\n期间所有消息被静默拦截',
]

txBox = slide.shapes.add_textbox(Inches(7.5), Inches(2.2), Inches(5), Inches(4))
tf = txBox.text_frame; tf.word_wrap = True
for j, lesson in enumerate(lessons):
    if j == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f'⚠️ {lesson}'
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'PingFang SC'
    p.space_before = Pt(10)
    p.space_after = Pt(4)

# ====== 第4页：行业穿透估值 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_GREEN; bar.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             '📊 Day2 行业穿透估值开发', font_size=30, color=DARK_TEXT, bold=True)

# 左侧：流程图
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5),
             '穿透估值原理', font_size=20, color=ACCENT_GREEN, bold=True)

steps = [
    ('①', 'ETF重仓股', '获取Top 10持仓明细', ACCENT_BLUE),
    ('②', '行业映射', '股票→申万一级行业\n(5854只股票/30行业)', ACCENT_GREEN),
    ('③', '权重计算', '按持仓占比加权\n(修复了%→小数bug)', ACCENT_ORANGE),
    ('④', '估值估算', '行业PE/PB × 权重\n= ETF穿透PE/PB分位', ACCENT_BLUE),
]

for i, (num, title, desc, color) in enumerate(steps):
    y = Inches(1.9 + i * 1.3)
    # 圆形编号
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.6), Inches(0.6))
    circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
    add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(0.6), Inches(0.5),
                 num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 内容
    add_text_box(slide, Inches(1.6), y, Inches(4.5), Inches(0.35),
                 title, font_size=18, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.35), Inches(4.5), Inches(0.7),
                 desc, font_size=13, color=SUBTITLE_GRAY)

# 右侧：数据卡片
card = add_shape(slide, Inches(7), Inches(1.2), Inches(5.8), Inches(5.5), LIGHT_GREEN, ACCENT_GREEN, Pt(2))
add_text_box(slide, Inches(7.3), Inches(1.4), Inches(5), Inches(0.5),
             '📈 覆盖率提升', font_size=20, color=ACCENT_GREEN, bold=True)

# 大数字
add_text_box(slide, Inches(7.3), Inches(2.0), Inches(2.5), Inches(1.2),
             '87.3%', font_size=52, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(9.8), Inches(2.4), Inches(2.5), Inches(0.5),
             'PE分位覆盖率', font_size=16, color=DARK_TEXT)
add_text_box(slide, Inches(9.8), Inches(2.9), Inches(2.5), Inches(0.5),
             '↑ 从67.9%提升19.4%', font_size=14, color=ACCENT_GREEN, bold=True)

# 分项数据
stats = [
    ('中证指数官方', '133只', '9.2%'),
    ('申万行业估值', '786只', '54.1%'),
    ('行业穿透估值', '282只', '19.4%'),
    ('巨潮资讯', '127只', '8.7%'),
    ('估算/降级', '36只', '2.5%'),
    ('无分位数据', '88只', '6.1%'),
]

txBox = slide.shapes.add_textbox(Inches(7.3), Inches(3.5), Inches(5.3), Inches(3))
tf = txBox.text_frame; tf.word_wrap = True
# 表头
for j, (name, count, pct) in enumerate(stats):
    if j == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    
    is_penetration = '穿透' in name
    color = ACCENT_GREEN if is_penetration else DARK_TEXT
    bold = is_penetration
    marker = '🆕' if is_penetration else '  '
    p.text = f'{marker} {name:<8} {count:>6}  {pct}'
    p.font.size = Pt(15)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'PingFang SC'
    p.space_before = Pt(6)
    p.space_after = Pt(2)

# ====== 第5页：数据来源与筛选结果 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_ORANGE; bar.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             '🎯 筛选结果变化', font_size=30, color=DARK_TEXT, bold=True)

# 三个对比卡片
comparisons = [
    ('正式低估池', '35只', '38只', '+3', ACCENT_GREEN, '满足PE/PB%≤30%\n+ PEG<1 + 日均≥1亿'),
    ('关注池', '5只', '106只', '+101', ACCENT_BLUE, '低估 + 日均≥300万\n穿透估值大幅补充'),
    ('数据不可用', '多', '0只', '大幅减少', ACCENT_ORANGE, '87.3%覆盖率\n仅88只QDII无数据'),
]

for i, (title, before, after, change, color, desc) in enumerate(comparisons):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(1.3), Inches(3.7), Inches(5.3), WHITE, color, Pt(2))
    
    add_text_box(slide, x + Inches(0.3), Inches(1.5), Inches(3.2), Inches(0.5),
                 title, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 优化前
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(1.3), Inches(0.4),
                 '优化前', font_size=13, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(2.7), Inches(1.3), Inches(0.7),
                 before, font_size=28, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 箭头
    add_text_box(slide, x + Inches(1.6), Inches(2.8), Inches(0.5), Inches(0.5),
                 '→', font_size=24, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    
    # 优化后
    add_text_box(slide, x + Inches(2.1), Inches(2.3), Inches(1.3), Inches(0.4),
                 '优化后', font_size=13, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(2.1), Inches(2.7), Inches(1.3), Inches(0.7),
                 after, font_size=28, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 变化标签
    change_shape = add_shape(slide, x + Inches(1.1), Inches(3.6), Inches(1.5), Inches(0.5), color)
    add_text_box(slide, x + Inches(1.1), Inches(3.65), Inches(1.5), Inches(0.4),
                 change, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 说明
    txBox = slide.shapes.add_textbox(x + Inches(0.3), Inches(4.4), Inches(3.2), Inches(2))
    tf = txBox.text_frame; tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = DARK_TEXT
            tf.paragraphs[0].font.name = 'PingFang SC'
        else:
            add_para(tf, line, font_size=14, color=SUBTITLE_GRAY)

# ====== 第6页：流水线架构 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             '🔄 流水线架构（优化后）', font_size=30, color=DARK_TEXT, bold=True)

pipeline_steps = [
    ('step0', '申万行业估值', '❌', False),
    ('step0a', '中证指数官方PE', '✅', False),
    ('step1', '行情采集', '✅', True),
    ('step2', '估值补全', '✅', True),
    ('step2a', '🆕穿透估值合并', '✅', False),
    ('step3', '低估筛选', '✅', True),
    ('step4', '结果对比', '✅', False),
    ('step5', '微信日报生成', '✅', True),
    ('step6', '指数快照归档', '✅', True),
    ('step7', '历史分位计算', '✅', False),
    ('step8', '趋势分析', '✅', False),
]

for i, (sid, name, status, required) in enumerate(pipeline_steps):
    x = Inches(0.5 + (i % 6) * 2.1)
    y = Inches(1.4 + (i // 6) * 2.6)
    
    is_new = '🆕' in name
    bg = LIGHT_GREEN if is_new else (LIGHT_BLUE if required else LIGHT_GRAY)
    border = ACCENT_GREEN if is_new else (ACCENT_BLUE if required else MED_GRAY)
    
    card = add_shape(slide, x, y, Inches(1.9), Inches(2.1), bg, border, Pt(2))
    
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(1.7), Inches(0.35),
                 sid, font_size=12, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.5), Inches(1.7), Inches(0.6),
                 name, font_size=15, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.2), Inches(1.7), Inches(0.35),
                 status, font_size=20, alignment=PP_ALIGN.CENTER)
    req_text = '必需' if required else '可选'
    req_color = ACCENT_BLUE if required else MED_GRAY
    add_text_box(slide, x + Inches(0.1), y + Inches(1.6), Inches(1.7), Inches(0.3),
                 req_text, font_size=12, color=req_color, alignment=PP_ALIGN.CENTER)

# 底部总结
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
             '⏱ 总耗时 ~9.5分钟 | 🆕 step2a穿透合并仅0.1秒 | 📅 日报推送09:40（确保读取当日数据）',
             font_size=15, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# ====== 第7页：下一步计划 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             '🚀 下一步优化方向', font_size=30, color=DARK_TEXT, bold=True)

plans = [
    ('P2', '关注池阈值优化', '关注池从5激增到106只\n需分级展示或调整筛选门槛', ACCENT_ORANGE, '中'),
    ('P3', '穿透估值定期更新', '每周手动运行持仓获取脚本\n保持穿透数据新鲜度', ACCENT_BLUE, '中'),
    ('P4', 'QDII/港股ETF估值', '88只无数据ETF主要是QDII\n考虑引入海外指数PE数据', MED_GRAY, '低'),
    ('P5', 'PEG数据补齐', 'PEG覆盖率0%，正式池PEG检查形同虚设\n需获取个股盈利增速数据', ACCENT_RED, '高'),
]

for i, (pid, title, desc, color, priority) in enumerate(plans):
    y = Inches(1.3 + i * 1.45)
    
    # 优先级标签
    pri_bg = ACCENT_RED if priority == '高' else (ACCENT_ORANGE if priority == '中' else MED_GRAY)
    pri_shape = add_shape(slide, Inches(0.8), y, Inches(0.9), Inches(1.1), pri_bg)
    add_text_box(slide, Inches(0.8), y + Inches(0.15), Inches(0.9), Inches(0.4),
                 pid, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), y + Inches(0.55), Inches(0.9), Inches(0.4),
                 priority, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 内容
    add_text_box(slide, Inches(1.9), y + Inches(0.05), Inches(5), Inches(0.4),
                 title, font_size=20, color=DARK_TEXT, bold=True)
    txBox = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.5), Inches(5), Inches(0.7))
    tf = txBox.text_frame; tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = SUBTITLE_GRAY
            tf.paragraphs[0].font.name = 'PingFang SC'
        else:
            add_para(tf, line, font_size=14, color=SUBTITLE_GRAY)

# ====== 保存 ======
out_path = Path('/Users/zhangxianghao/.qclaw/workspace/etf-agent/output/ETF优化汇报_20260515.pptx')
prs.save(str(out_path))
print(f'✅ PPT已保存: {out_path}')
print(f'   文件大小: {out_path.stat().st_size / 1024:.1f}KB')
